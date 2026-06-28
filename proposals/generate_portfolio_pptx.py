#!/usr/bin/env python3
"""Generate Alison Granger Portfolio PPTX in Blooming brand style."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Colors (white background, bosque accents) ──
BG_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
BOSQUE = RGBColor(0x5C, 0x6B, 0x5A)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_60 = RGBColor(0x6B, 0x6B, 0x6B)
TIERRA_LIGHT = RGBColor(0x8B, 0x7D, 0x6B)
CREAM = DARK
CREAM_60 = DARK_60

SLIDE_W = Inches(11)
SLIDE_H = Inches(8.5)


def set_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_text_box(slide, left, top, width, height, text, font_name="Calibri",
                 font_size=12, color=CREAM, bold=False, italic=False,
                 alignment=PP_ALIGN.LEFT, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox, tf


def add_paragraph(tf, text, font_name="Calibri", font_size=12, color=CREAM,
                  bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                  space_before=0, line_spacing=1.15):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    p.line_spacing = Pt(font_size * line_spacing)
    return p


def add_footer(slide, left_text="Alison Granger \u00b7 Portfolio", right_text=""):
    y = Inches(7.7)
    add_text_box(slide, Inches(1.2), y, Inches(3), Inches(0.3),
                 left_text, font_size=7, color=CREAM_60)
    if right_text:
        add_text_box(slide, Inches(7.8), y, Inches(2), Inches(0.3),
                     right_text, font_size=7, color=CREAM_60,
                     alignment=PP_ALIGN.RIGHT)


def add_line(slide, left, top, width=Inches(0.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(0.75))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
    shape.line.fill.background()
    return shape


def add_label(slide, left, top, text):
    add_text_box(slide, left, top, Inches(4), Inches(0.3),
                 text.upper(), font_size=8, color=TIERRA_LIGHT, bold=False)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    L = Inches(1.2)
    W = Inches(8.6)

    # ═══════════════════════════════════════
    # SLIDE 1: COVER
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)

    add_text_box(slide, L, Inches(1.2), Inches(3), Inches(0.4),
                 "BLOOMING", font_size=10, color=CREAM)

    add_text_box(slide, L, Inches(3.0), Inches(8), Inches(1.6),
                 "Alison Granger\nPortfolio",
                 font_name="Georgia", font_size=42, color=CREAM,
                 line_spacing=1.1)

    add_text_box(slide, L, Inches(5.0), Inches(6), Inches(0.4),
                 "Event strategist & ecosystem architect \u00b7 10+ years of operating experience",
                 font_size=11, color=CREAM_60)

    add_text_box(slide, L, Inches(5.7), Inches(2), Inches(0.35),
                 "CONFIDENTIAL", font_size=7, color=TIERRA_LIGHT)

    add_footer(slide, left_text="Blooming", right_text="weareblooming.co")

    # ═══════════════════════════════════════
    # SLIDE 2: POSITIONING
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)

    add_label(slide, L, Inches(0.9), "Positioning")

    add_text_box(slide, L, Inches(1.5), W, Inches(1.0),
                 "Most events optimize for attendance.\nAlison optimizes for outcomes.",
                 font_name="Georgia", font_size=28, color=CREAM, line_spacing=1.15)

    add_text_box(slide, L, Inches(2.8), Inches(7.5), Inches(1.0),
                 "Designing curated environments where the right people meet, build, and move the ecosystem forward. 10+ years bridging product, GTM, and developer relations through narrative-driven ecosystem events across AI, developer, and startup ecosystems.",
                 font_size=10, color=CREAM_60, line_spacing=1.7)

    stats = [("150+", "Events architected"), ("15k+", "Builders engaged"),
             ("10+", "Years operating"), ("SF", "Based \u00b7 Global reach")]
    for i, (num, label) in enumerate(stats):
        x = Inches(1.2 + i * 2.2)
        add_text_box(slide, x, Inches(4.3), Inches(2), Inches(0.6),
                     num, font_name="Georgia", font_size=28, color=CREAM)
        add_text_box(slide, x, Inches(4.95), Inches(2), Inches(0.3),
                     label.upper(), font_size=7, color=TIERRA_LIGHT)

    add_footer(slide, right_text="02")

    # ═══════════════════════════════════════
    # SLIDE 3: POINT OF VIEW
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)

    add_label(slide, L, Inches(0.9), "Point of view")

    add_text_box(slide, L, Inches(1.5), W, Inches(0.6),
                 "Why most events fail",
                 font_name="Georgia", font_size=28, color=CREAM)

    cols = [
        ("The problem", "Attendance over impact",
         "They optimize for scale when they should optimize for signal. Attendees become a metric, not a curated ecosystem."),
        ("What creates value", "Intentional curation",
         "Proximity to decision-makers and pioneers. Access to unfiltered patterns. Communities built on repeated, high-trust interactions."),
        ("Why now", "AI moves too fast for playbooks",
         "Companies expanding cross-border need ecosystem architecture, not logistics. Someone who understands developer communities, international dynamics, and what moves the needle."),
    ]
    for i, (label, title, desc) in enumerate(cols):
        x = Inches(1.2 + i * 2.9)
        y = Inches(2.6)
        add_text_box(slide, x, y, Inches(2.5), Inches(0.25),
                     label, font_name="Georgia", font_size=10, color=TIERRA_LIGHT)
        add_line(slide, x, y + Inches(0.35))
        add_text_box(slide, x, y + Inches(0.5), Inches(2.5), Inches(0.4),
                     title, font_name="Georgia", font_size=18, color=CREAM)
        add_text_box(slide, x, y + Inches(1.0), Inches(2.5), Inches(1.2),
                     desc, font_size=9.5, color=CREAM_60, line_spacing=1.65)

    add_footer(slide, right_text="03")

    # ═══════════════════════════════════════
    # SLIDE 4: OPERATING LAYER
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)

    add_label(slide, L, Inches(0.9), "Operating layer")

    add_text_box(slide, L, Inches(1.5), W, Inches(0.6),
                 "Product \u00d7 GTM \u00d7 DevRel",
                 font_name="Georgia", font_size=28, color=CREAM)

    add_text_box(slide, L, Inches(2.4), Inches(7.5), Inches(1.5),
                 "Across all work\u2014from European festivals and productions to AI ecosystem programs\u2014Alison operates at the intersection of Product, GTM, and Developer Relations. She architects narrative arcs that transform fragmented corporate initiatives into continuous engagement loops. Rather than executing logistics, she designs systems where every developer touchpoint reinforces positioning and builds momentum.",
                 font_size=10, color=CREAM_60, line_spacing=1.7)

    add_text_box(slide, L, Inches(4.5), Inches(8), Inches(0.4),
                 "Narrative design  \u00b7  Ecosystem strategy  \u00b7  Community architecture  \u00b7  Real-world execution",
                 font_name="Georgia", font_size=16, color=TIERRA_LIGHT, italic=True)

    add_footer(slide, right_text="04")

    # ═══════════════════════════════════════
    # SLIDE 5: SELECTED EXPERIENCE
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)

    add_label(slide, L, Inches(0.9), "Selected experience \u00b7 2022\u20132026")

    add_text_box(slide, L, Inches(1.5), W, Inches(0.6),
                 "Building ecosystems at scale",
                 font_name="Georgia", font_size=28, color=CREAM)

    exp = [
        ("01", "AI User Group", "Early operator & builder",
         "Scaled AIUG from launch through monthly meetups, a 5-city AI engineering workshop tour, and 2 large conferences in SF. From pre-event strategy to post-event storytelling: partner communication, audience engagement, social media amplification, livestream production, and community narratives that extend beyond the room."),
        ("02", "Enterprise AI Activations", "Strategic operator \u00b7 2023\u20132026",
         "Operated large-scale ecosystem programs with NVIDIA, Google, IBM, Meta. Designed multi-event series as sustained narrative arcs. Connected product launches, developer education, and community engagement into continuous loops."),
        ("03", "Blooming", "Founder & CEO \u00b7 2026",
         "Consolidation of 10+ years of operating experience into a strategic layer. Launched with 5 ecosystem events in 60 days across US\u2013LatAm corridors\u2014hackathons, developer meetups, and side events\u2014testing and validating Blooming\u2019s model of high-signal curation and ecosystem design."),
    ]
    for i, (num, title, timing, desc) in enumerate(exp):
        x = Inches(1.2 + i * 2.9)
        y = Inches(2.6)
        add_text_box(slide, x, y, Inches(2.5), Inches(0.25),
                     num, font_name="Georgia", font_size=10, color=TIERRA_LIGHT)
        add_line(slide, x, y + Inches(0.35))
        add_text_box(slide, x, y + Inches(0.5), Inches(2.5), Inches(0.4),
                     title, font_name="Georgia", font_size=18, color=CREAM)
        add_text_box(slide, x, y + Inches(1.0), Inches(2.5), Inches(0.25),
                     timing.upper(), font_size=7, color=TIERRA_LIGHT)
        add_text_box(slide, x, y + Inches(1.35), Inches(2.5), Inches(1.8),
                     desc, font_size=9.5, color=CREAM_60, line_spacing=1.65)

    add_footer(slide, right_text="05")

    # ═══════════════════════════════════════
    # SLIDE 6: PARTNERS
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)

    add_label(slide, L, Inches(0.9), "Partners")

    add_text_box(slide, L, Inches(1.5), W, Inches(0.6),
                 "Trusted by teams building the future",
                 font_name="Georgia", font_size=28, color=CREAM)

    add_text_box(slide, L, Inches(2.4), Inches(7.5), Inches(0.6),
                 "Strategic event operations and ecosystem programs delivered in partnership with leading technology companies.",
                 font_size=10, color=CREAM_60, line_spacing=1.7)

    partners = ["Google", "NVIDIA", "IBM", "Miro", "Intercom", "Canva",
                "DataDog", "n8n", "UiPath", "Pipelex", "TRAE", "The AI Collective"]

    # Two rows of partner names
    row1 = partners[:6]
    row2 = partners[6:]
    for i, name in enumerate(row1):
        x = Inches(1.2 + i * 1.5)
        add_text_box(slide, x, Inches(3.5), Inches(1.4), Inches(0.4),
                     name, font_size=12, color=CREAM_60, bold=False)
    for i, name in enumerate(row2):
        x = Inches(1.2 + i * 1.5)
        add_text_box(slide, x, Inches(4.2), Inches(1.4), Inches(0.4),
                     name, font_size=12, color=CREAM_60, bold=False)

    add_text_box(slide, L, Inches(5.3), Inches(7.5), Inches(0.4),
                 "From developer workshops and conferences to multi-event ecosystem programs across AI and infrastructure.",
                 font_size=9, color=CREAM_60, italic=True)

    add_footer(slide, right_text="06")

    # ═══════════════════════════════════════
    # SLIDE 7: CORE SERVICES
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)

    add_label(slide, L, Inches(0.9), "Core services")

    add_text_box(slide, L, Inches(1.5), W, Inches(0.6),
                 "What Blooming delivers",
                 font_name="Georgia", font_size=28, color=CREAM)

    services = [
        ("Event Strategy & Positioning", [
            "Multi-event series aligned to business objectives",
            "Narrative architecture that builds momentum",
            "Audience mapping and community segmentation",
        ]),
        ("Community Architecture", [
            "Ecosystem design for new markets and geographies",
            "Community lead recruitment and training",
            "Curation and signal-to-noise optimization",
        ]),
        ("Full Execution", [
            "End-to-end operations and logistics",
            "On-site production and community management",
            "Attendee curation and experience design",
        ]),
        ("Ecosystem Programs", [
            "Developer programs and partnerships",
            "Sponsorship strategy and partner alignment",
            "International expansion programs",
        ]),
    ]
    for i, (title, items) in enumerate(services):
        col = i % 2
        row = i // 2
        x = Inches(1.2 + col * 4.5)
        y = Inches(2.6 + row * 2.2)
        add_text_box(slide, x, y, Inches(3.8), Inches(0.35),
                     title, font_name="Georgia", font_size=15, color=TIERRA_LIGHT, italic=True)
        for j, item in enumerate(items):
            add_text_box(slide, x, y + Inches(0.45 + j * 0.35), Inches(3.8), Inches(0.3),
                         f"\u2013  {item}", font_size=9.5, color=CREAM_60)

    add_footer(slide, right_text="07")

    # ═══════════════════════════════════════
    # SLIDE 8: SIGNATURE STRENGTHS
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)

    add_label(slide, L, Inches(0.9), "Why Alison")

    add_text_box(slide, L, Inches(1.5), W, Inches(0.6),
                 "Signature strengths",
                 font_name="Georgia", font_size=28, color=CREAM)

    strengths = [
        ("01", "Operational velocity",
         "Moves fast, ships clean. No quality trade-offs. 150+ events = a detailed map of what works and what doesn\u2019t."),
        ("02", "Curation over scale",
         "Builds rooms, not crowds. Every attendee adds signal. Understands what moves developers and founders. Thinks in networks."),
        ("03", "Cross-border expertise",
         "Operates across US, Europe, and international geographies with local understanding. Systems thinking that treats events as long-term ecosystem infrastructure."),
    ]
    for i, (num, title, desc) in enumerate(strengths):
        x = Inches(1.2 + i * 2.9)
        y = Inches(2.6)
        add_text_box(slide, x, y, Inches(2.5), Inches(0.25),
                     num, font_name="Georgia", font_size=10, color=TIERRA_LIGHT)
        add_line(slide, x, y + Inches(0.35))
        add_text_box(slide, x, y + Inches(0.5), Inches(2.5), Inches(0.4),
                     title, font_name="Georgia", font_size=18, color=CREAM)
        add_text_box(slide, x, y + Inches(1.0), Inches(2.5), Inches(1.2),
                     desc, font_size=9.5, color=CREAM_60, line_spacing=1.65)

    add_text_box(slide, L, Inches(5.5), Inches(7), Inches(0.6),
                 "Blooming is built from hands-on operating experience, not consulting theory. It combines executional excellence with narrative design\u2014treating events not as logistics, but as infrastructure for ecosystem growth.",
                 font_name="Georgia", font_size=12, color=TIERRA_LIGHT, italic=True, line_spacing=1.6)

    add_footer(slide, right_text="08")

    # ═══════════════════════════════════════
    # SLIDE 9: ABOUT BLOOMING
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)

    add_label(slide, L, Inches(0.9), "About Blooming")

    add_text_box(slide, L, Inches(1.5), W, Inches(1.0),
                 "A strategic operating layer for\ncross-border growth",
                 font_name="Georgia", font_size=28, color=CREAM, line_spacing=1.15)

    # Two columns
    add_text_box(slide, L, Inches(3.0), Inches(3.8), Inches(0.35),
                 "What it is", font_name="Georgia", font_size=18, color=CREAM)
    add_text_box(slide, L, Inches(3.5), Inches(3.8), Inches(1.5),
                 "A strategic operating layer for tech companies expanding across markets, communities, and narratives. It sits between ecosystem strategy, community design, and real-world execution\u2014turning events into structured moments of growth, positioning, and signal.",
                 font_size=9.5, color=CREAM_60, line_spacing=1.65)

    add_text_box(slide, Inches(5.8), Inches(3.0), Inches(3.8), Inches(0.35),
                 "Who it\u2019s for", font_name="Georgia", font_size=18, color=CREAM)
    add_text_box(slide, Inches(5.8), Inches(3.5), Inches(3.8), Inches(1.5),
                 "Series A+ and growth-stage companies operating across borders\u2014especially in AI and developer ecosystems\u2014where traditional events no longer create differentiation or meaningful engagement.",
                 font_size=9.5, color=CREAM_60, line_spacing=1.65)

    # Bottom details
    add_text_box(slide, L, Inches(5.5), Inches(3), Inches(0.35),
                 "SF-based", font_name="Georgia", font_size=14, color=TIERRA_LIGHT, italic=True)
    add_text_box(slide, L, Inches(5.9), Inches(3), Inches(0.3),
                 "WITH GLOBAL REACH", font_size=7, color=TIERRA_LIGHT)

    add_text_box(slide, Inches(4.5), Inches(5.5), Inches(4), Inches(0.35),
                 "Co-founders", font_name="Georgia", font_size=14, color=TIERRA_LIGHT, italic=True)
    add_text_box(slide, Inches(4.5), Inches(5.9), Inches(4), Inches(0.3),
                 "ANDR\u00c9S (GTM), NACHO (GROWTH), CLAUDIO (DESIGN)", font_size=7, color=TIERRA_LIGHT)

    add_footer(slide, right_text="09")

    # ═══════════════════════════════════════
    # SLIDE 10: CLOSE
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)

    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(8), Inches(1.0),
                 "Let\u2019s build\nsomething together.",
                 font_name="Georgia", font_size=38, color=TIERRA_LIGHT,
                 italic=True, alignment=PP_ALIGN.CENTER, line_spacing=1.1)

    add_text_box(slide, Inches(2.5), Inches(3.8), Inches(6), Inches(1.0),
                 "The best ecosystems aren\u2019t planned from a distance\u2014they\u2019re built by someone in the room who understands the stakes, the people, and the signal.",
                 font_name="Georgia", font_size=12, color=CREAM_60,
                 italic=True, alignment=PP_ALIGN.CENTER, line_spacing=1.7)

    add_line(slide, Inches(5.0), Inches(5.0), Inches(1))

    add_text_box(slide, Inches(2.5), Inches(5.3), Inches(6), Inches(0.35),
                 "alison@weareblooming.co",
                 font_size=13, color=CREAM_60, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2.5), Inches(5.7), Inches(6), Inches(0.3),
                 "+33 (0) 6 48 15 53 98",
                 font_size=9, color=CREAM_60, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2.5), Inches(6.1), Inches(6), Inches(0.3),
                 "WEAREBLOOMING.CO",
                 font_size=8, color=TIERRA_LIGHT, alignment=PP_ALIGN.CENTER)

    add_footer(slide, left_text="Blooming", right_text="weareblooming.co")

    # ── Save ──
    output = "/Users/panthervillagran/blooming-website/proposals/alison-granger-portfolio.pptx"
    prs.save(output)
    print(f"Saved: {output}")


if __name__ == "__main__":
    main()
