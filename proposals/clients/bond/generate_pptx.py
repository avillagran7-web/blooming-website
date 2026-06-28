from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ───────────────────────────────────────────────
NEGRO   = RGBColor(0x1A, 0x1A, 0x1A)
BLANCO  = RGBColor(0xFF, 0xFF, 0xFF)
BOSQUE  = RGBColor(0x5C, 0x6B, 0x5A)
TIERRA  = RGBColor(0x8B, 0x7D, 0x6B)   # logo horizontal ellipses
TIERRA_L= RGBColor(0xC4, 0xB8, 0xA8)   # light tierra / borders
BORDER  = RGBColor(0xE0, 0xDD, 0xD8)
GRAY    = RGBColor(0x55, 0x55, 0x55)
GRAY_L  = RGBColor(0x88, 0x88, 0x88)
GRAY_F  = RGBColor(0xAA, 0xA8, 0xA4)

# ── Slide size 16:9 ──────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

CG = "Cormorant Garamond"
SG = "Space Grotesk"

# ── Helpers ───────────────────────────────────────────────
def bg(slide):
    add_rect(slide, 0, 0, W, H, fill_color=BLANCO)

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill_color:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color; shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_oval(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(9, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_line(slide, x1, y1, x2, y2, color=BORDER, width=Pt(0.5)):
    shape = slide.shapes.add_connector(1, x1, y1, x2, y2)
    shape.line.color.rgb = color; shape.line.width = width
    return shape

def add_textbox(slide, text, x, y, w, h, font_name=SG, font_size=11,
                bold=False, italic=False, color=NEGRO, align=PP_ALIGN.LEFT, word_wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = word_wrap
    tf = tb.text_frame; tf.word_wrap = word_wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name; run.font.size = Pt(font_size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_multiline(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, word_wrap=True):
    """lines = [(text, font_name, font_size, bold, italic, color, space_before_pt)]"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = word_wrap
    tf = tb.text_frame; tf.word_wrap = word_wrap
    first = True
    for (text, fname, fsize, fbold, fitalic, fcolor, spb) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if spb: p.space_before = Pt(spb)
        run = p.add_run()
        run.text = text
        run.font.name = fname; run.font.size = Pt(fsize)
        run.font.bold = fbold; run.font.italic = fitalic
        run.font.color.rgb = fcolor
    return tb

def accent_line(slide, x=Inches(0.85), y=Inches(1.08)):
    add_line(slide, x, y, x + Inches(0.32), y, color=NEGRO, width=Pt(0.75))

def label(slide, text, x=Inches(0.85), y=Inches(0.82)):
    add_textbox(slide, text.upper(), x, y, Inches(10), Inches(0.22),
                font_name=SG, font_size=7, color=BOSQUE)

def footer(slide, left="Blooming × Bond.ai", right="weareblooming.co"):
    y = H - Inches(0.48)
    add_textbox(slide, left,  Inches(0.85), y, Inches(5), Inches(0.28),
                font_name=SG, font_size=6, color=GRAY_F)
    add_textbox(slide, right, Inches(8.0),  y, Inches(5), Inches(0.28),
                font_name=SG, font_size=6, color=GRAY_F, align=PP_ALIGN.RIGHT)

def isotipo(slide, cx, cy, scale=1.0):
    """Blooming isotipo — proper ellipses, proper brand colors"""
    s = scale
    # Vertical ellipses (bosque)
    add_oval(slide, cx - Inches(0.065*s), cy - Inches(0.26*s), Inches(0.13*s), Inches(0.32*s), BOSQUE)
    add_oval(slide, cx - Inches(0.065*s), cy - Inches(0.06*s), Inches(0.13*s), Inches(0.32*s), BOSQUE)
    # Horizontal ellipses (tierra)
    add_oval(slide, cx - Inches(0.34*s), cy - Inches(0.065*s), Inches(0.32*s), Inches(0.13*s), TIERRA)
    add_oval(slide, cx + Inches(0.02*s),  cy - Inches(0.065*s), Inches(0.32*s), Inches(0.13*s), TIERRA)
    # Center dot
    r = Inches(0.065*s)
    add_oval(slide, cx - r, cy - r, r*2, r*2, NEGRO)

# ═══════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK); bg(s1)

isotipo(s1, Inches(1.08), Inches(0.72), scale=0.9)
add_textbox(s1, "BLOOMING", Inches(1.38), Inches(0.62), Inches(3), Inches(0.28),
            font_name=SG, font_size=8, color=NEGRO)

add_multiline(s1, Inches(0.85), Inches(1.5), Inches(10), Inches(2.5), [
    ("Bond.ai Community × Blooming", CG, 42, False, False, NEGRO, 0),
    ("Pilot Partnership.",           CG, 42, False, True,  BOSQUE, 4),
])

add_textbox(s1, "AI Hackathon  ·  San Francisco  ·  May 2026",
            Inches(0.85), Inches(3.85), Inches(8), Inches(0.32),
            font_name=SG, font_size=9, color=GRAY_L)

# Date badge
add_rect(s1, Inches(0.85), Inches(4.38), Inches(1.3), Inches(0.26), line_color=BORDER)
add_textbox(s1, "APRIL 2026", Inches(0.92), Inches(4.41), Inches(1.15), Inches(0.2),
            font_name=SG, font_size=6, color=GRAY_L)

add_textbox(s1, "Prepared by Alison Granger, Co-Founder & CEO",
            Inches(0.85), Inches(6.5), Inches(7), Inches(0.28),
            font_name=SG, font_size=8, color=GRAY_L)

add_multiline(s1, Inches(9.5), Inches(6.5), Inches(3.5), Inches(0.28), [
    ("BLOOMING  ×  BOND.AI", SG, 7, False, False, GRAY_L, 0),
], align=PP_ALIGN.RIGHT)

footer(s1, "Blooming", "weareblooming.co")

# ═══════════════════════════════════════════════════════════
# SLIDE 2 — THE OPPORTUNITY
# ═══════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK); bg(s2)
accent_line(s2); label(s2, "The Opportunity")

add_multiline(s2, Inches(0.85), Inches(1.2), Inches(10), Inches(1.3), [
    ("Two operators.", CG, 30, False, False, NEGRO, 0),
    ("One gap worth closing.", CG, 30, False, True, BOSQUE, 4),
])

add_rect(s2, Inches(0.85), Inches(2.62), Inches(0.04), Inches(0.72), fill_color=TIERRA_L)
add_textbox(s2,
    "This Tuesday, Sahar and I met for the first time, which was a great encounter, considering how many "
    "events we've both been at on the same Luma page. We've been orbiting the same ecosystem for several years. "
    "This is good momentum.",
    Inches(1.05), Inches(2.6), Inches(10.5), Inches(0.8),
    font_name=CG, font_size=10.5, italic=True, color=GRAY)

col_y = Inches(3.55)
for cx, title, items in [
    (Inches(0.85), "What Bond brings", [
        "Audience trust — companies follow AI Tidbits, attend Bond events",
        "Sponsor relationships built over years in the SF AI ecosystem",
        "Inbound demand from companies who can't replicate that reach",
    ]),
    (Inches(7.0), "What Blooming brings", [
        "Production engine — 150+ events: IBM, NVIDIA, Google, Canva",
        "Alison Granger AI Ecosystem network — 70K+ reach, 15k dev community credibility in SF",
        "Flawless execution under tight timelines, every time, each timezone.",
    ]),
]:
    add_textbox(s2, title, cx, col_y, Inches(5.8), Inches(0.28),
                font_name=CG, font_size=13, italic=True, color=BOSQUE)
    add_multiline(s2, cx, col_y + Inches(0.35), Inches(5.8), Inches(1.3), [
        (f"–  {it}", SG, 8.5, False, False, GRAY, 0 if i == 0 else 4)
        for i, it in enumerate(items)
    ])

add_line(s2, Inches(0.85), Inches(5.3), Inches(12.5), Inches(5.3), color=BORDER)
add_textbox(s2,
    "Sahar brings audience reach and sponsor relationships. Alison brings production excellence and developer "
    "community credibility. Together, we deliver what companies we know expect at the level they demand.",
    Inches(0.85), Inches(5.42), Inches(11.5), Inches(0.7),
    font_name=CG, font_size=10, italic=True, color=GRAY_L)

footer(s2)

# ═══════════════════════════════════════════════════════════
# SLIDE 3 — THE PILOT
# ═══════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK); bg(s3)
accent_line(s3); label(s3, "The Pilot")

add_multiline(s3, Inches(0.85), Inches(1.2), Inches(10), Inches(1.2), [
    ("One event to", CG, 30, False, False, NEGRO, 0),
    ("test the format.", CG, 30, False, True, BOSQUE, 4),
])

specs = [
    ("Format",    "1-Day AI Hackathon"),
    ("Attendees", "80 Qualified Builders"),
    ("Location",  "San Francisco"),
    ("Timing",    "Mid–End of May 2026"),
    ("Sponsors",  "Starting at 2 Companies"),
]
sx = Inches(0.85)
for lbl_t, val_t in specs:
    add_textbox(s3, lbl_t.upper(), sx, Inches(2.48), Inches(2.4), Inches(0.18),
                font_name=SG, font_size=6, color=BOSQUE)
    add_textbox(s3, val_t, sx, Inches(2.68), Inches(2.4), Inches(0.32),
                font_name=CG, font_size=12, color=NEGRO)
    sx += Inches(2.45)

add_line(s3, Inches(0.85), Inches(3.1), Inches(12.5), Inches(3.1), color=BORDER)

for cx, title, items in [
    (Inches(0.85), "Bond · AI Tidbits", [
        "Leads sponsor outreach and closes deals",
        "Pays event expenses upfront to Blooming Events.",
        "Brings Bond audience and AI Tidbits community",
    ]),
    (Inches(7.0), "Blooming Events by Alison Granger", [
        "Leads full production and logistics online",
        "Brings curated attendees list for participant outreach",
        "Manages day-of presence the day of the hackathon",
    ]),
]:
    add_textbox(s3, title, cx, Inches(3.22), Inches(5.8), Inches(0.28),
                font_name=CG, font_size=13, italic=True, color=BOSQUE)
    add_multiline(s3, cx, Inches(3.58), Inches(5.8), Inches(1.2), [
        (f"–  {it}", SG, 8.5, False, False, GRAY, 0 if i == 0 else 4)
        for i, it in enumerate(items)
    ])

add_line(s3, Inches(0.85), Inches(4.98), Inches(12.5), Inches(4.98), color=BORDER)

for num, lbl_t, mx in [
    ("80%+", "Attendee satisfaction",  Inches(0.85)),
    ("100+", "Events architected",     Inches(4.85)),
    ("10+",  "Years operating",        Inches(8.85)),
]:
    add_textbox(s3, num, mx, Inches(5.12), Inches(3.5), Inches(0.52),
                font_name=CG, font_size=24, color=NEGRO)
    add_textbox(s3, lbl_t.upper(), mx, Inches(5.68), Inches(3.5), Inches(0.22),
                font_name=SG, font_size=6, color=BOSQUE)

footer(s3)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 — HACKATHON THEMES
# ═══════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK); bg(s4)
accent_line(s4); label(s4, "Hackathon Themes")

add_multiline(s4, Inches(0.85), Inches(1.2), Inches(10), Inches(1.0), [
    ("Three", CG, 30, False, False, NEGRO, 0),
    ("directions.", CG, 30, False, True, BOSQUE, 0),
])

themes = [
    ("AI Agents in Production",       "Build agents that ship real impact through demos."),
    ("AI for Developer Productivity",  "Build AI tools that make developers 10× more productive."),
    ("AI for GTM & Growth",            "Build AI tools for sales, marketing, and growth teams."),
]
ty = Inches(2.35)
for i, (title, subtitle) in enumerate(themes):
    add_line(s4, Inches(0.85), ty, Inches(12.5), ty, color=BORDER)
    ty += Inches(0.1)
    add_textbox(s4, title,    Inches(0.85), ty, Inches(11), Inches(0.32),
                font_name=CG, font_size=16, color=NEGRO)
    add_textbox(s4, subtitle, Inches(0.85), ty + Inches(0.35), Inches(11), Inches(0.25),
                font_name=SG, font_size=8, italic=True, color=GRAY_L)
    ty += Inches(1.0)

add_line(s4, Inches(0.85), ty, Inches(12.5), ty, color=BORDER)

add_rect(s4, Inches(0.85), ty + Inches(0.12), Inches(0.03), Inches(0.55), fill_color=TIERRA_L)
add_textbox(s4,
    "Work in progress — this is the core work of the collaboration. "
    "Example: \"My ICP is primarily startup founders (Series A and below) — "
    "do you have a sense of the demographics of the attendees?\" · 15k attendees list · "
    "Portfolio for reference · Last past event: high interest in Live Demo format · 120 attendees",
    Inches(1.05), ty + Inches(0.1), Inches(11), Inches(0.7),
    font_name=CG, font_size=9, italic=True, color=GRAY_L)

footer(s4)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 — IMMEDIATE DELIVERABLES
# ═══════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK); bg(s5)
accent_line(s5); label(s5, "Once Contract Signed")

add_multiline(s5, Inches(0.85), Inches(1.2), Inches(10), Inches(1.0), [
    ("Immediate", CG, 30, False, True, BOSQUE, 0),
    ("deliverables.", CG, 30, False, False, NEGRO, 0),
])

deliverables = [
    ("01", "Overall Schedule", [
        "Sharing trusted inputs about pre-event, day-of, and post-event milestones",
    ]),
    ("02", "Operations Plan", [
        "Pre-event coordination — venue, sponsors, judges, mentors, platforms",
        "Day-of execution — setup, programming, logistics",
        "Post-event wrap — deliverables, follow-ups, reporting",
    ]),
    ("03", "Work Session: Define the Hackathon", [
        "Final theme selection — Agents, Dev Productivity, or GTM",
        "Challenge structure and prerequisites",
        "Launch publication and messaging",
        "Sponsor outreach strategy",
    ]),
]
cx = Inches(0.85)
cw = Inches(3.9)
for num, title, items in deliverables:
    add_textbox(s5, num, cx, Inches(2.45), cw, Inches(0.22),
                font_name=CG, font_size=10, color=BOSQUE)
    add_textbox(s5, title, cx, Inches(2.7), cw, Inches(0.38),
                font_name=CG, font_size=14, color=NEGRO)
    add_multiline(s5, cx, Inches(3.18), cw, Inches(2.5), [
        (f"–  {it}", SG, 8.5, False, False, GRAY, 0 if i == 0 else 5)
        for i, it in enumerate(items)
    ])
    cx += Inches(4.22)

footer(s5)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 — SPONSOR PACKAGE
# ═══════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK); bg(s6)
accent_line(s6); label(s6, "Sponsorship Package — Starting at $2k · Pilot test · Aiming for $5k")

add_multiline(s6, Inches(0.85), Inches(1.2), Inches(10), Inches(1.2), [
    ("80 qualified builders.", CG, 30, False, False, NEGRO, 0),
    ("8 hours. One room.",     CG, 30, False, True,  BOSQUE, 4),
])

pkg_cols = [
    ("Visibility", [
        "Equal logo on all materials — website, emails, signage, excellent & unique atmosphere",
        "Event co-branding: 'AI Hackathon powered by [Sponsors]'",
        "LinkedIn, newsletter & X mentions",
    ]),
    ("Engagement", [
        "10-minute speaking slot — demo, API walkthrough, or use case",
        "Booth space with good placement",
        "Optional: 'Best use of [Company] API' prize ($500–$1,500)",
        "API credits distributed to all participants",
        "Dedicated Discord channel",
    ]),
    ("Data & Access", [
        "Opt-in attendee list — names, emails, LinkedIn",
        "Warm intros to winning teams",
        "Post-event report — photos, metrics, feedback",
        "Recruiting pipeline — access to top talent in the room",
    ]),
]
cx = Inches(0.85)
for title, items in pkg_cols:
    add_textbox(s6, title, cx, Inches(2.6), Inches(4.0), Inches(0.28),
                font_name=CG, font_size=13, italic=True, color=BOSQUE)
    add_multiline(s6, cx, Inches(2.96), Inches(4.0), Inches(2.4), [
        (f"–  {it}", SG, 8, False, False, GRAY, 0 if i == 0 else 4)
        for i, it in enumerate(items)
    ])
    cx += Inches(4.22)

add_line(s6, Inches(0.85), Inches(5.55), Inches(12.5), Inches(5.55), color=BORDER)
add_textbox(s6,
    "Prizes are paid directly by sponsors outside the sponsor fee. "
    "Each sponsor defines their own prize amount, format, and category — cash, API credits, or product licenses.",
    Inches(0.85), Inches(5.67), Inches(11.5), Inches(0.6),
    font_name=CG, font_size=10, italic=True, color=GRAY_L)

footer(s6)

# ═══════════════════════════════════════════════════════════
# SLIDE 7 — SCOPE OF WORK
# ═══════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK); bg(s7)
accent_line(s7); label(s7, "Blooming — Scope of Work")

add_multiline(s7, Inches(0.85), Inches(1.2), Inches(10), Inches(1.0), [
    ("The", CG, 30, False, False, NEGRO, 0),
    ("Event Engine.", CG, 30, False, True, BOSQUE, 0),
])

scope_data = [
    ("Pre-Event", [
        "Venue sourcing, negotiation & contract",
        "Catering coordination (80 attendees)",
        "AV/tech setup & recording",
        "Secondary sponsor outreach via Alison network",
        "Sponsor coordination — challenges, prizes",
        "Curated database outreach & registration",
        "Devpost setup for submissions & judging",
        "Discord server — participant channels",
    ]),
    ("Day-Of", [
        "7am: Setup — tables, chairs, signage, AV",
        "Registration desk — check-in, badges, swag",
        "Emcee/host coordination",
        "Speaker & workshop timing management",
        "Mentor circulation coordination",
        "Real-time troubleshooting",
        "Demo & judging logistics (5pm–7pm)",
        "Winner announcements & prize handoff",
        "Venue teardown & equipment return",
    ]),
    ("Post-Event", [
        "Event report — attendance, photos, folder, metrics",
        "Opt-in attendee list to sponsors",
        "Intro emails — sponsors to winning teams",
        "Thank you emails to all participants",
        "Post-event survey & feedback collection",
    ]),
]
cx = Inches(0.85)
for title, items in scope_data:
    add_textbox(s7, title, cx, Inches(2.32), Inches(4.0), Inches(0.28),
                font_name=CG, font_size=13, color=NEGRO)
    add_line(s7, cx, Inches(2.65), cx + Inches(3.8), Inches(2.65), color=BORDER)
    add_multiline(s7, cx, Inches(2.75), Inches(4.0), Inches(2.55), [
        (f"–  {it}", SG, 7.5, False, False, GRAY, 0 if i == 0 else 3)
        for i, it in enumerate(items)
    ])
    cx += Inches(4.22)

# Extras row
ey = Inches(5.52)
add_line(s7, Inches(0.85), ey, Inches(12.5), ey, color=BORDER)
ey += Inches(0.12)
for ex_label, ex_items in [
    ("Team Onsite", [
        "1 lead operator",
        "Volunteers: Morning (Set Up) · Launch (Catering) · Evening (Final Demo)",
        "Curation of trusted partners: F&B, AV, Venues",
    ]),
    ("Scope to be defined", ["Judges & mentors recruitment & briefing"]),
    ("Add-On Deliverables", ["Social media content — LinkedIn & X", "Videos & photos"]),
]:
    add_textbox(s7, ex_label.upper(), Inches(0.85) if ex_label=="Team Onsite" else
                (Inches(5.5) if "Scope" in ex_label else Inches(9.0)),
                ey, Inches(4.0), Inches(0.18), font_name=SG, font_size=6, color=BOSQUE)
    add_multiline(s7,
                  Inches(0.85) if ex_label=="Team Onsite" else
                  (Inches(5.5) if "Scope" in ex_label else Inches(9.0)),
                  ey + Inches(0.2), Inches(4.0), Inches(0.65), [
        (f"–  {it}", SG, 7.5, False, False, GRAY_L, 0 if i == 0 else 3)
        for i, it in enumerate(ex_items)
    ])

footer(s7)

# ═══════════════════════════════════════════════════════════
# SLIDE 8 — INVESTMENT & REVENUE
# ═══════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK); bg(s8)
accent_line(s8); label(s8, "Investment & Revenue")

add_multiline(s8, Inches(0.85), Inches(1.2), Inches(10), Inches(1.0), [
    ("The Pilot.", CG, 30, False, False, NEGRO, 0),
    ("Let's make it work.", CG, 30, False, True, BOSQUE, 4),
])

# Blooming fee
add_rect(s8, Inches(0.85), Inches(2.35), Inches(11.65), Inches(0.38), line_color=BORDER)
add_textbox(s8, "BLOOMING FLAT FEE", Inches(1.0), Inches(2.42), Inches(2.5), Inches(0.22),
            font_name=SG, font_size=6.5, color=BOSQUE)
add_textbox(s8, "$1,000", Inches(3.2), Inches(2.38), Inches(2), Inches(0.32),
            font_name=CG, font_size=17, color=NEGRO)
add_textbox(s8, "Covers Alison pre-event management, online day-of production support & post-event deliverables.",
            Inches(5.0), Inches(2.43), Inches(7.2), Inches(0.28),
            font_name=SG, font_size=7.5, color=GRAY_L)

scenarios = [
    ("Starting Point — $2–2.5k / sponsor · 4 sponsors", [
        ("Venue (free — cleaning fee)", "$300–$500"),
        ("Venue (paid)",                "$1,500–$5,000"),
        ("Catering (80 ppl)",           "$2,500–$3,000"),
        ("Contingency (10%)",           "$250–$900"),
        ("On-site production team",     "$1,225"),
    ], "Revenue (4 sponsors)", "$8,000–$10,000", [
        ("Best case (free venue)", "~$3,000"),
        ("Standard (paid venue)",  "Up to $2,000"),
        ("Break-even",             "2–4 sponsors"),
    ]),
    ("Expected Scenario — $5k / sponsor · 3–5 sponsors", [
        ("Venue (free — cleaning fee)", "$300–$500"),
        ("Venue (paid)",                "$3,500–$6,000"),
        ("Catering (80 ppl)",           "$3,000–$5,000"),
        ("Contingency (10%)",           "$650–$1,100"),
        ("On-site production team",     "$2,000"),
    ], "Revenue (3–5 sponsors)", "$15,000–$25,000", [
        ("Best case (free venue)", "$7,960–$17,960"),
        ("Standard (paid venue)",  "$3,175–$13,175"),
        ("Break-even",             "2–3 sponsors"),
    ]),
]

for i, (title, rows, rev_label, revenue, profit_rows) in enumerate(scenarios):
    cx = Inches(0.85) if i == 0 else Inches(7.1)
    cw = Inches(5.8)
    add_textbox(s8, title, cx, Inches(2.9), cw, Inches(0.25),
                font_name=CG, font_size=11, italic=True, color=BOSQUE)
    ry = Inches(3.22)
    for item, price in rows:
        add_line(s8, cx, ry, cx + cw, ry, color=BORDER)
        ry += Inches(0.06)
        add_textbox(s8, item,  cx,              ry, Inches(3.8), Inches(0.24), font_name=SG, font_size=7.5, color=GRAY)
        add_textbox(s8, price, cx + Inches(3.9), ry, Inches(1.8), Inches(0.24), font_name=CG, font_size=9.5, color=NEGRO, align=PP_ALIGN.RIGHT)
        ry += Inches(0.27)
    add_line(s8, cx, ry, cx + cw, ry, color=NEGRO, width=Pt(0.75))
    ry += Inches(0.06)
    add_textbox(s8, rev_label, cx,              ry, Inches(3.8), Inches(0.3), font_name=SG, font_size=8.5, bold=True, color=NEGRO)
    add_textbox(s8, revenue,   cx + Inches(3.9), ry, Inches(1.8), Inches(0.3), font_name=CG, font_size=13, color=BOSQUE, align=PP_ALIGN.RIGHT)
    ry += Inches(0.38)
    add_textbox(s8, "PROFIT", cx, ry, Inches(2), Inches(0.18), font_name=SG, font_size=6, color=BOSQUE)
    ry += Inches(0.22)
    for pname, pval in profit_rows:
        add_line(s8, cx, ry, cx + cw, ry, color=BORDER)
        ry += Inches(0.05)
        add_textbox(s8, pname, cx,              ry, Inches(3.8), Inches(0.22), font_name=SG, font_size=7.5, color=GRAY)
        add_textbox(s8, pval,  cx + Inches(3.9), ry, Inches(1.8), Inches(0.22), font_name=CG, font_size=9.5, color=NEGRO, align=PP_ALIGN.RIGHT)
        ry += Inches(0.25)

# Profit split
py = Inches(6.28)
add_line(s8, Inches(0.85), py, Inches(12.5), py, color=BORDER)
py += Inches(0.1)
add_textbox(s8, "PROFIT SPLIT", Inches(0.85), py, Inches(1.8), Inches(0.18), font_name=SG, font_size=6, color=BOSQUE)
add_textbox(s8, "80%", Inches(2.85), py - Inches(0.02), Inches(1.2), Inches(0.32), font_name=CG, font_size=18, color=NEGRO)
add_textbox(s8, "Bond", Inches(2.85), py + Inches(0.3), Inches(1.2), Inches(0.18), font_name=SG, font_size=6, color=BOSQUE)
add_textbox(s8, "20%", Inches(4.25), py - Inches(0.02), Inches(1.2), Inches(0.32), font_name=CG, font_size=18, color=NEGRO)
add_textbox(s8, "Blooming", Inches(4.25), py + Inches(0.3), Inches(1.5), Inches(0.18), font_name=SG, font_size=6, color=BOSQUE)
add_textbox(s8, "Media production: price on demand  ·  $1.5k–$3k pocket",
            Inches(6.5), py + Inches(0.05), Inches(6.0), Inches(0.22),
            font_name=CG, font_size=9.5, italic=True, color=GRAY_L, align=PP_ALIGN.RIGHT)

footer(s8)

# ═══════════════════════════════════════════════════════════
# SLIDE 9 — CLOSE
# ═══════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK); bg(s9)

isotipo(s9, W/2, Inches(1.3), scale=1.1)

add_multiline(s9, Inches(1.5), Inches(1.85), Inches(10.3), Inches(1.9), [
    ("1 pilot event.", CG, 38, False, False, NEGRO, 0),
    ("Crash test the format.", CG, 38, False, True, BOSQUE, 4),
], align=PP_ALIGN.CENTER)

add_textbox(s9, "Let's ship what we do best.",
            Inches(1.5), Inches(3.72), Inches(10.3), Inches(0.32),
            font_name=CG, font_size=13, italic=True, color=GRAY_L, align=PP_ALIGN.CENTER)

add_textbox(s9,
    "The ideal way to validate a partnership is to do the work. "
    "One hackathon in SF — ideally mid or end of May. Date to be confirmed upon validation.",
    Inches(2.2), Inches(4.12), Inches(8.9), Inches(0.7),
    font_name=CG, font_size=11.5, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

steps = [
    "Align on hackathon theme — pick the one with the strongest sponsor momentum",
    "Bond leads sponsor outreach, Blooming starts venue sourcing in parallel",
    "Lock 2 sponsors to confirm the pilot — go from there",
]
sy = Inches(4.95)
for step in steps:
    add_oval(s9, Inches(4.28), sy + Inches(0.08), Inches(0.07), Inches(0.07), BOSQUE)
    add_textbox(s9, step, Inches(4.45), sy, Inches(4.5), Inches(0.28),
                font_name=SG, font_size=8.5, color=GRAY)
    sy += Inches(0.34)

add_textbox(s9, "alison@weareblooming.co",
            Inches(1.5), Inches(6.1), Inches(10.3), Inches(0.28),
            font_name=SG, font_size=10, color=GRAY, align=PP_ALIGN.CENTER)
add_textbox(s9, "WEAREBLOOMING.CO",
            Inches(1.5), Inches(6.4), Inches(10.3), Inches(0.22),
            font_name=SG, font_size=7, color=BOSQUE, align=PP_ALIGN.CENTER)

footer(s9, "Blooming", "Growth that matters")

# ── Save ─────────────────────────────────────────────────
out = "proposals/clients/bond/Blooming x Bond.ai — Pilot Partnership.pptx"
prs.save(out)
print(f"Saved → {out}")
