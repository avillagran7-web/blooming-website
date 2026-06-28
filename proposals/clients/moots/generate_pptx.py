#!/usr/bin/env python3
"""Generate Blooming x Moots proposal PPTX from the HTML deck content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors (white background version) ──
BG_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
BOSQUE = RGBColor(0x5C, 0x6B, 0x5A)
BOSQUE_DEEP = RGBColor(0x4A, 0x57, 0x48)
DARK = RGBColor(0x1A, 0x1A, 0x1A)        # primary text on white
DARK_60 = RGBColor(0x6B, 0x6B, 0x6B)     # secondary text
TIERRA_LIGHT = RGBColor(0x8B, 0x7D, 0x6B) # accent
TIERRA = RGBColor(0x8B, 0x7D, 0x6B)
# Aliases for readability
CREAM = DARK          # headings
CREAM_60 = DARK_60    # body text

# Slide dimensions: 11 x 8.5 inches (landscape letter)
SLIDE_W = Inches(11)
SLIDE_H = Inches(8.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_text_box(slide, left, top, width, height, text, font_name="Calibri",
                 font_size=12, color=CREAM, bold=False, italic=False,
                 alignment=PP_ALIGN.LEFT, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox, tf


def add_paragraph(tf, text, font_name="Calibri", font_size=12, color=CREAM,
                  bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                  space_before=0, line_spacing=1.15):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    p.line_spacing = Pt(font_size * line_spacing)
    return p


def add_footer(slide, left_text="Blooming x Moots", right_text=""):
    y = Inches(7.7)
    add_text_box(slide, Inches(1.2), y, Inches(3), Inches(0.3),
                 left_text, font_size=7, color=CREAM_60)
    if right_text:
        add_text_box(slide, Inches(7.8), y, Inches(2), Inches(0.3),
                     right_text, font_size=7, color=CREAM_60,
                     alignment=PP_ALIGN.RIGHT)


def add_line(slide, left, top, width=Inches(0.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(0.75))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
    shape.line.fill.background()
    return shape


def add_label(slide, left, top, text):
    """Small uppercase label."""
    add_text_box(slide, left, top, Inches(4), Inches(0.3),
                 text.upper(), font_size=8, color=TIERRA_LIGHT, bold=False)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank

    L = Inches(1.2)  # left margin
    W = Inches(8.6)  # content width

    # ═══════════════════════════════════════
    # SLIDE 1: COVER
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BOSQUE)

    add_text_box(slide, L, Inches(1.2), Inches(3), Inches(0.4),
                 "BLOOMING", font_size=10, color=CREAM, bold=False)

    add_text_box(slide, L, Inches(3.0), Inches(8), Inches(1.6),
                 "Bringing Moots to the\nAI ecosystem.",
                 font_name="Georgia", font_size=42, color=CREAM,
                 italic=False, line_spacing=1.1)

    add_text_box(slide, L, Inches(5.0), Inches(6), Inches(0.4),
                 "Go-to-Market Partnership  \u00b7  May \u2013 July 2026",
                 font_size=11, color=CREAM_60)

    add_text_box(slide, L, Inches(5.7), Inches(2), Inches(0.35),
                 "CONFIDENTIAL", font_size=7, color=TIERRA_LIGHT)

    add_footer(slide, right_text="weareblooming.co")

    # ═══════════════════════════════════════
    # SLIDE 2: THE OPPORTUNITY
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BOSQUE)

    add_label(slide, L, Inches(0.9), "The opportunity")

    add_text_box(slide, L, Inches(1.5), W, Inches(1.0),
                 "Events are the last trust channel.\nMost companies waste them.",
                 font_name="Georgia", font_size=28, color=CREAM, line_spacing=1.15)

    add_text_box(slide, L, Inches(2.8), Inches(7.5), Inches(1.2),
                 "\"Tonje, I saw Moots at HumanX \u2014 and I immediately understood what you\u2019re building. As someone who has produced over fifty events, I know the pain firsthand: guest selection is manual, follow-ups fall through the cracks, and no one can tell you which conversation actually moved the needle. Moots solves exactly this. The timing couldn\u2019t be better.\"",
                 font_name="Georgia", font_size=12, color=TIERRA_LIGHT,
                 italic=True, line_spacing=1.7)

    add_text_box(slide, L, Inches(4.5), Inches(7.5), Inches(0.9),
                 "AI is saturating every digital channel. Cold outreach is noise. But events remain one of the last high-signal ways to build real trust \u2014 yet most organizations treat them as one-off moments with no memory. Guest curation is reactive, attribution from event to pipeline is nonexistent, and every gathering starts from zero.",
                 font_size=10, color=CREAM_60, line_spacing=1.7)

    # Three highlights
    highlights = [
        ("The shift", "Events are becoming the primary trust channel as AI saturates digital outreach. High-signal, high-touch, irreplaceable."),
        ("The gap", "Most organizations still treat events as one-off moments \u2014 no memory, no compounding, no attribution."),
        ("The timing", "Moots is ready. The audience exists. Field marketing teams are actively looking for this solution."),
    ]
    for i, (title, desc) in enumerate(highlights):
        x = Inches(1.2 + i * 2.9)
        y = Inches(5.8)
        add_text_box(slide, x, y, Inches(2.5), Inches(0.35),
                     title, font_name="Georgia", font_size=14, color=TIERRA_LIGHT, italic=True)
        add_line(slide, x, Inches(6.2))
        add_text_box(slide, x, Inches(6.35), Inches(2.5), Inches(0.8),
                     desc, font_size=9, color=CREAM_60, line_spacing=1.6)

    add_footer(slide, right_text="02")

    # ═══════════════════════════════════════
    # SLIDE 3: WHAT WE BRING
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BOSQUE)

    add_label(slide, L, Inches(0.9), "What we bring")

    add_text_box(slide, L, Inches(1.5), W, Inches(1.0),
                 "Not a referral partner.\nA live proof of concept.",
                 font_name="Georgia", font_size=30, color=CREAM, line_spacing=1.15)

    add_text_box(slide, L, Inches(2.8), Inches(7), Inches(1.0),
                 "Blooming produces curated events for AI and technology companies across three ecosystems. We don\u2019t just make introductions \u2014 we put Moots in front of the right people through warm relationships, curated programming, and a live environment where the platform speaks for itself. Leaders from IBM, Nvidia, Neo4j, and Google already sit in our rooms.",
                 font_size=10, color=CREAM_60, line_spacing=1.7)

    # Traction numbers
    traction = [("50+", "Events produced"), ("70K+", "Attendees reached"),
                ("30+", "Warm contacts"), ("3", "Ecosystems")]
    for i, (num, label) in enumerate(traction):
        x = Inches(1.2 + i * 2.2)
        add_text_box(slide, x, Inches(4.3), Inches(2), Inches(0.6),
                     num, font_name="Georgia", font_size=28, color=CREAM)
        add_text_box(slide, x, Inches(4.95), Inches(2), Inches(0.3),
                     label.upper(), font_size=7, color=TIERRA_LIGHT)

    # Values
    add_text_box(slide, L, Inches(5.7), Inches(6), Inches(0.4),
                 "Speed  \u00b7  Quality  \u00b7  Trust",
                 font_name="Georgia", font_size=16, color=TIERRA_LIGHT, italic=True)

    add_footer(slide, right_text="03")

    # ═══════════════════════════════════════
    # SLIDE 4: THE PLAN
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BOSQUE)

    add_label(slide, L, Inches(0.9), "The plan")

    add_text_box(slide, L, Inches(1.5), W, Inches(0.6),
                 "Three phases. One clear path.",
                 font_name="Georgia", font_size=28, color=CREAM)

    phases = [
        ("01", "Awareness", "May 2026",
         "Introduce Moots across Blooming\u2019s network through warm, personal outreach. Position the platform with the right people before any pitch \u2014 context first, always."),
        ("02", "Demos", "May \u2013 June 2026",
         "Curated one-on-one demos with field marketing managers and event leads who run events constantly. Warm introductions, not cold lists."),
        ("03", "The Event", "June \u2013 July 2026",
         "A hands-on workshop powered by Moots \u2014 attendees use the platform live, share best practices, and give real-time feedback. The product becomes the proof."),
    ]
    for i, (num, title, timing, desc) in enumerate(phases):
        x = Inches(1.2 + i * 2.9)
        y_base = Inches(2.8)
        add_text_box(slide, x, y_base, Inches(2.5), Inches(0.25),
                     num, font_name="Georgia", font_size=10, color=TIERRA_LIGHT)
        add_line(slide, x, y_base + Inches(0.35))
        add_text_box(slide, x, y_base + Inches(0.5), Inches(2.5), Inches(0.4),
                     title, font_name="Georgia", font_size=18, color=CREAM)
        add_text_box(slide, x, y_base + Inches(1.0), Inches(2.5), Inches(0.25),
                     timing.upper(), font_size=7, color=TIERRA_LIGHT)
        add_text_box(slide, x, y_base + Inches(1.35), Inches(2.5), Inches(1.2),
                     desc, font_size=9.5, color=CREAM_60, line_spacing=1.65)

    add_footer(slide, right_text="04")

    # ═══════════════════════════════════════
    # SLIDE 5: THE EVENT
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BOSQUE)

    add_label(slide, L, Inches(0.9), "The event")

    add_text_box(slide, L, Inches(1.5), W, Inches(1.0),
                 "A hands-on workshop,\npowered by Moots.",
                 font_name="Georgia", font_size=28, color=CREAM, line_spacing=1.15)

    add_text_box(slide, L, Inches(2.8), Inches(7.5), Inches(1.0),
                 "An intimate session with 10\u201315 field marketing managers from leading AI and tech companies. Everyone gets hands on the platform \u2014 using Moots live, sharing best practices for guest intelligence, and giving real-time feedback. They don\u2019t just hear about the product. They use it, together.",
                 font_size=10, color=CREAM_60, line_spacing=1.7)

    details = [("10\u201315", "Guests"), ("San Francisco", "Location"),
               ("2.5\u20133 hrs", "Duration"), ("Field Marketing", "Audience")]
    for i, (val, label) in enumerate(details):
        x = Inches(1.2 + i * 2.2)
        add_text_box(slide, x, Inches(4.3), Inches(2), Inches(0.5),
                     val, font_name="Georgia", font_size=18, color=CREAM)
        add_text_box(slide, x, Inches(4.85), Inches(2), Inches(0.3),
                     label.upper(), font_size=7, color=TIERRA_LIGHT)

    add_text_box(slide, L, Inches(5.6), Inches(8), Inches(0.4),
                 "The 6th event on Moots produces 7x more pipeline than the 1st \u2014 because the platform compounds intelligence.",
                 font_name="Georgia", font_size=12, color=TIERRA_LIGHT, italic=True)

    add_footer(slide, right_text="05")

    # ═══════════════════════════════════════
    # SLIDE 6: WHY THIS WORKS
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BOSQUE)

    add_label(slide, L, Inches(0.9), "Why this works")

    add_text_box(slide, L, Inches(1.5), W, Inches(1.0),
                 "Not just distribution.\nA growth engine.",
                 font_name="Georgia", font_size=28, color=CREAM, line_spacing=1.15)

    reasons = [
        ("01", "Warm network",
         "Real relationships, not cold lists. Every introduction comes with context, trust, and a reason to take the meeting. Built over fifty events and three ecosystems."),
        ("02", "Aligned incentives",
         "Commission-based pricing means Blooming is invested in Moots\u2019 success \u2014 not just delivery. We win when you convert, which means every introduction is intentional."),
        ("03", "The right people, at the right moment",
         "Field marketing managers who run events constantly. They feel the pain of manual guest selection and broken attribution every week. They\u2019re ready for Moots."),
    ]
    for i, (num, title, desc) in enumerate(reasons):
        x = Inches(1.2 + i * 2.9)
        y_base = Inches(2.8)
        add_text_box(slide, x, y_base, Inches(2.5), Inches(0.25),
                     num, font_name="Georgia", font_size=10, color=TIERRA_LIGHT)
        add_line(slide, x, y_base + Inches(0.35))
        add_text_box(slide, x, y_base + Inches(0.5), Inches(2.5), Inches(0.4),
                     title, font_name="Georgia", font_size=18, color=CREAM)
        add_text_box(slide, x, y_base + Inches(1.0), Inches(2.5), Inches(1.2),
                     desc, font_size=9.5, color=CREAM_60, line_spacing=1.65)

    add_footer(slide, right_text="06")

    # ═══════════════════════════════════════
    # SLIDE 7: INVESTMENT
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BOSQUE)

    add_label(slide, L, Inches(0.9), "Investment")

    add_text_box(slide, L, Inches(1.5), W, Inches(0.6),
                 "Simple, transparent, aligned.",
                 font_name="Georgia", font_size=28, color=CREAM)

    add_text_box(slide, L, Inches(2.4), Inches(4), Inches(0.7),
                 "$3,000", font_name="Georgia", font_size=36, color=CREAM)

    add_text_box(slide, L, Inches(3.1), Inches(5), Inches(0.3),
                 "starting at \u2014 excl. logistics & commission",
                 font_size=8, color=CREAM_60)

    items = [
        ("Awareness campaign (warm intros across network)", "$1,000"),
        ("Curated demos (per qualified meeting)", "$50 / each"),
        ("Curated workshop (production & curation)", "$1,000"),
        ("Event logistics (venue, F&B)", "At cost"),
        ("Revenue commission (6-month window)", "10%"),
    ]

    y = Inches(3.7)
    add_line(slide, L, y, Inches(7.5))
    for i, (item, price) in enumerate(items):
        row_y = y + Inches(0.15 + i * 0.45)
        add_text_box(slide, L, row_y, Inches(5.5), Inches(0.35),
                     item, font_size=10, color=CREAM_60)
        add_text_box(slide, Inches(7.5), row_y, Inches(2), Inches(0.35),
                     price, font_name="Georgia", font_size=13, color=CREAM,
                     alignment=PP_ALIGN.RIGHT)

    add_text_box(slide, L, Inches(6.3), Inches(6.5), Inches(0.5),
                 "Pricing reflects early-partnership rates for a first activation. Market references available upon request.",
                 font_name="Georgia", font_size=10, color=TIERRA_LIGHT, italic=True,
                 line_spacing=1.6)

    add_footer(slide, right_text="07")

    # ═══════════════════════════════════════
    # SLIDE 8: TIMELINE
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BOSQUE)

    add_label(slide, L, Inches(0.9), "Timeline")

    add_text_box(slide, L, Inches(1.5), W, Inches(0.6),
                 "Before May 1 \u2192 End of July",
                 font_name="Georgia", font_size=28, color=CREAM)

    timeline = [
        ("Before May 1", "Align on messaging, define target profiles, finalize partnership terms"),
        ("May 1 \u2013 15", "Launch warm outreach \u2014 introduce Moots across Blooming\u2019s network with personal context"),
        ("May 15 \u2013 June 15", "Run curated demos with qualified field marketing leads \u2014 one-on-one, warm introductions"),
        ("June \u2013 Early July", "Hands-on workshop in San Francisco \u2014 powered by Moots, 10\u201315 field marketing managers"),
        ("July", "Follow-up, pipeline tracking, outcomes review \u2014 set the foundation for ongoing activation"),
    ]

    for i, (date, desc) in enumerate(timeline):
        row_y = Inches(2.6 + i * 0.75)
        add_text_box(slide, L, row_y, Inches(1.8), Inches(0.35),
                     date.upper(), font_size=7, color=TIERRA_LIGHT)
        add_text_box(slide, Inches(3.2), row_y, Inches(6), Inches(0.5),
                     desc, font_size=10, color=CREAM_60, line_spacing=1.6)
        if i < len(timeline) - 1:
            add_line(slide, L, row_y + Inches(0.5), Inches(8))

    add_footer(slide, right_text="08")

    # ═══════════════════════════════════════
    # SLIDE 9: CLOSE
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BOSQUE)

    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(8), Inches(1.0),
                 "Let\u2019s build the right room.",
                 font_name="Georgia", font_size=38, color=TIERRA_LIGHT,
                 italic=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2.5), Inches(3.8), Inches(6), Inches(1.4),
                 "\"Tonje, we\u2019re both French, we both ended up in the middle of the AI world, and we both believe relationships are what actually move things forward. I\u2019d love to help you get Moots in front of the people who need it most \u2014 and I think we can do something genuinely great together.\"",
                 font_name="Georgia", font_size=12, color=CREAM_60,
                 italic=True, alignment=PP_ALIGN.CENTER, line_spacing=1.7)

    add_line(slide, Inches(5.0), Inches(5.4), Inches(1))

    add_text_box(slide, Inches(2.5), Inches(5.7), Inches(6), Inches(0.35),
                 "alison@weareblooming.co",
                 font_size=13, color=CREAM_60, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2.5), Inches(6.1), Inches(6), Inches(0.3),
                 "WEAREBLOOMING.CO",
                 font_size=8, color=TIERRA_LIGHT, alignment=PP_ALIGN.CENTER)

    add_footer(slide, right_text="Confidential")

    # ── Save ──
    output = "/Users/panthervillagran/blooming-website/proposals/clients/moots/blooming-x-moots-proposal.pptx"
    prs.save(output)
    print(f"Saved: {output}")


if __name__ == "__main__":
    main()
