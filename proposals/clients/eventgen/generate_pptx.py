from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ───────────────────────────────────────────────
NEGRO   = RGBColor(0x1A, 0x1A, 0x1A)
BLANCO  = RGBColor(0xFF, 0xFF, 0xFF)
BOSQUE  = RGBColor(0x5C, 0x6B, 0x5A)
TIERRA  = RGBColor(0x8B, 0x7D, 0x6B)
TIERRA_L= RGBColor(0xC4, 0xB8, 0xA8)
BORDER  = RGBColor(0xE0, 0xDD, 0xD8)
GRAY    = RGBColor(0x55, 0x55, 0x55)
GRAY_L  = RGBColor(0x88, 0x88, 0x88)
GRAY_F  = RGBColor(0xAA, 0xA8, 0xA4)

# ── Slide size 16:9 ──────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

CG = "Cormorant Garamond"
SG = "Space Grotesk"

# ── Helpers ───────────────────────────────────────────────
def bg(slide):
    add_rect(slide, 0, 0, W, H, fill_color=BLANCO)

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill_color:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color; shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_oval(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(9, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_line(slide, x1, y1, x2, y2, color=BORDER, width=Pt(0.5)):
    shape = slide.shapes.add_connector(1, x1, y1, x2, y2)
    shape.line.color.rgb = color; shape.line.width = width
    return shape

def add_textbox(slide, text, x, y, w, h, font_name=SG, font_size=11,
                bold=False, italic=False, color=NEGRO, align=PP_ALIGN.LEFT, word_wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = word_wrap
    tf = tb.text_frame; tf.word_wrap = word_wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name; run.font.size = Pt(font_size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_multiline(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, word_wrap=True):
    """lines = [(text, font_name, font_size, bold, italic, color, space_before_pt)]"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = word_wrap
    tf = tb.text_frame; tf.word_wrap = word_wrap
    first = True
    for (text, fname, fsize, fbold, fitalic, fcolor, spb) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if spb: p.space_before = Pt(spb)
        run = p.add_run()
        run.text = text
        run.font.name = fname; run.font.size = Pt(fsize)
        run.font.bold = fbold; run.font.italic = fitalic
        run.font.color.rgb = fcolor
    return tb

def accent_line(slide, x=Inches(0.85), y=Inches(1.08)):
    add_line(slide, x, y, x + Inches(0.32), y, color=NEGRO, width=Pt(0.75))

def label(slide, text, x=Inches(0.85), y=Inches(0.82)):
    add_textbox(slide, text.upper(), x, y, Inches(10), Inches(0.22),
                font_name=SG, font_size=7, color=BOSQUE)

def footer(slide, left="Blooming × EventGen", right="weareblooming.co"):
    y = H - Inches(0.48)
    add_textbox(slide, left,  Inches(0.85), y, Inches(5), Inches(0.28),
                font_name=SG, font_size=6, color=GRAY_F)
    add_textbox(slide, right, Inches(8.0),  y, Inches(5), Inches(0.28),
                font_name=SG, font_size=6, color=GRAY_F, align=PP_ALIGN.RIGHT)

def isotipo(slide, cx, cy, scale=1.0):
    s = scale
    add_oval(slide, cx - Inches(0.065*s), cy - Inches(0.26*s), Inches(0.13*s), Inches(0.32*s), BOSQUE)
    add_oval(slide, cx - Inches(0.065*s), cy - Inches(0.06*s), Inches(0.13*s), Inches(0.32*s), BOSQUE)
    add_oval(slide, cx - Inches(0.34*s),  cy - Inches(0.065*s), Inches(0.32*s), Inches(0.13*s), TIERRA)
    add_oval(slide, cx + Inches(0.02*s),  cy - Inches(0.065*s), Inches(0.32*s), Inches(0.13*s), TIERRA)
    r = Inches(0.065*s)
    add_oval(slide, cx - r, cy - r, r*2, r*2, NEGRO)

# ═══════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK); bg(s1)

isotipo(s1, Inches(1.08), Inches(0.72), scale=0.9)
add_textbox(s1, "BLOOMING", Inches(1.38), Inches(0.62), Inches(3), Inches(0.28),
            font_name=SG, font_size=8, color=NEGRO)

add_multiline(s1, Inches(0.85), Inches(1.5), Inches(10), Inches(2.5), [
    ("EventGen × Blooming", CG, 42, False, False, NEGRO, 0),
    ("Venues & Vendors Sourcing.",  CG, 42, False, True,  BOSQUE, 4),
])

add_textbox(s1, "Multi-City Event Infrastructure  ·  10+ Cities  ·  2-Week Delivery",
            Inches(0.85), Inches(3.85), Inches(9), Inches(0.32),
            font_name=SG, font_size=9, color=GRAY_L)

add_rect(s1, Inches(0.85), Inches(4.38), Inches(1.3), Inches(0.26), line_color=BORDER)
add_textbox(s1, "MAY 2026", Inches(0.92), Inches(4.41), Inches(1.15), Inches(0.2),
            font_name=SG, font_size=6, color=GRAY_L)

add_textbox(s1, "Prepared by Alison Granger, Co-Founder & CEO",
            Inches(0.85), Inches(6.5), Inches(7), Inches(0.28),
            font_name=SG, font_size=8, color=GRAY_L)

add_multiline(s1, Inches(9.5), Inches(6.5), Inches(3.5), Inches(0.28), [
    ("BLOOMING  ×  EVENTGEN", SG, 7, False, False, GRAY_L, 0),
], align=PP_ALIGN.RIGHT)

footer(s1, "Blooming", "weareblooming.co")

# ═══════════════════════════════════════════════════════════
# SLIDE 2 — THE OPPORTUNITY
# ═══════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK); bg(s2)
accent_line(s2); label(s2, "The Opportunity")

add_multiline(s2, Inches(0.85), Inches(1.2), Inches(10), Inches(1.3), [
    ("One platform.", CG, 30, False, False, NEGRO, 0),
    ("Every city. Every format.", CG, 30, False, True, BOSQUE, 4),
])

add_rect(s2, Inches(0.85), Inches(2.62), Inches(0.04), Inches(0.72), fill_color=TIERRA_L)
add_textbox(s2,
    "EventGen is an integrated booking system combining venue selection, catering, photography/video, "
    "printing, and badge systems — all synced with Google Calendar. The goal: organize any event in ~5 minutes, "
    "consistently, across 10+ cities globally.",
    Inches(1.05), Inches(2.6), Inches(10.5), Inches(0.8),
    font_name=CG, font_size=10.5, italic=True, color=GRAY)

col_y = Inches(3.55)
for cx, title, items in [
    (Inches(0.85), "What EventGen brings", [
        "Integrated booking engine — venue, catering, AV, badges in one platform",
        "Google Calendar sync for real-time availability across all cities",
        "Scalable model: Meetup · Workshop · Hackathon — up to 200 attendees",
    ]),
    (Inches(7.0), "What Blooming brings", [
        "Alison knows Digital Jungle and its stakeholders — she was part of the launch via BeMyApp / AI User Group",
        "Network in SF, LatAm, and Paris tech ecosystems — sourcing at speed",
        "Continuity in the collaboration means execution from day one.",
    ]),
]:
    add_textbox(s2, title, cx, col_y, Inches(5.8), Inches(0.28),
                font_name=CG, font_size=13, italic=True, color=BOSQUE)
    add_multiline(s2, cx, col_y + Inches(0.35), Inches(5.8), Inches(1.3), [
        (f"–  {it}", SG, 8.5, False, False, GRAY, 0 if i == 0 else 4)
        for i, it in enumerate(items)
    ])

add_line(s2, Inches(0.85), Inches(5.3), Inches(12.5), Inches(5.3), color=BORDER)
add_textbox(s2,
    "EventGen needs a trusted sourcing partner with real relationships in each city. "
    "Blooming has the network, the execution track record, and the continuity with Digital Jungle to move fast.",
    Inches(0.85), Inches(5.42), Inches(11.5), Inches(0.7),
    font_name=CG, font_size=10, italic=True, color=GRAY_L)

footer(s2)

# ═══════════════════════════════════════════════════════════
# SLIDE 3 — THE MISSION
# ═══════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK); bg(s3)
accent_line(s3); label(s3, "The Mission")

add_multiline(s3, Inches(0.85), Inches(1.2), Inches(10), Inches(1.2), [
    ("Build the supplier ecosystem.", CG, 30, False, False, NEGRO, 0),
    ("2 weeks.", CG, 30, False, True, BOSQUE, 4),
])

specs = [
    ("Timeline",    "2 Weeks — from May 16"),
    ("Pilot City",  "San Francisco / Digital Jungle"),
    ("Cities",      "10 Cities Globally"),
    ("Formats",     "Meetup · Workshop · Hackathon"),
    ("Revenue",     "20% Commission Model"),
]
sx = Inches(0.85)
for lbl_t, val_t in specs:
    add_textbox(s3, lbl_t.upper(), sx, Inches(2.48), Inches(2.4), Inches(0.18),
                font_name=SG, font_size=6, color=BOSQUE)
    add_textbox(s3, val_t, sx, Inches(2.68), Inches(2.4), Inches(0.32),
                font_name=CG, font_size=11, color=NEGRO)
    sx += Inches(2.45)

add_line(s3, Inches(0.85), Inches(3.1), Inches(12.5), Inches(3.1), color=BORDER)

for cx, title, items in [
    (Inches(0.85), "City Targets", [
        "London · Paris · Berlin",
        "San Francisco · New York · Los Angeles · Austin · Seattle",
        "Toronto · Singapore (TBD — timezone friction)",
    ]),
    (Inches(7.0), "Event Format Requirements", [
        "Meetup — 4pm to 9pm",
        "Workshop — 9am to 5pm",
        "Hackathon — 1 day (8am–9pm) or 2 days",
        "Up to 150–200 attendees per format",
    ]),
]:
    add_textbox(s3, title, cx, Inches(3.22), Inches(5.8), Inches(0.28),
                font_name=CG, font_size=13, italic=True, color=BOSQUE)
    add_multiline(s3, cx, Inches(3.58), Inches(5.8), Inches(1.2), [
        (f"–  {it}", SG, 8.5, False, False, GRAY, 0 if i == 0 else 4)
        for i, it in enumerate(items)
    ])

add_line(s3, Inches(0.85), Inches(4.98), Inches(12.5), Inches(4.98), color=BORDER)

for num, lbl_t, mx in [
    ("10+",  "Cities globally",       Inches(0.85)),
    ("3",    "Event formats",          Inches(4.85)),
    ("2wk",  "Delivery timeline",      Inches(8.85)),
]:
    add_textbox(s3, num, mx, Inches(5.12), Inches(3.5), Inches(0.52),
                font_name=CG, font_size=24, color=NEGRO)
    add_textbox(s3, lbl_t.upper(), mx, Inches(5.68), Inches(3.5), Inches(0.22),
                font_name=SG, font_size=6, color=BOSQUE)

footer(s3)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 — SCOPE OF WORK
# ═══════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK); bg(s4)
accent_line(s4); label(s4, "Scope of Work")

add_multiline(s4, Inches(0.85), Inches(1.2), Inches(10), Inches(1.0), [
    ("Three tracks.", CG, 30, False, False, NEGRO, 0),
    ("One deadline.", CG, 30, False, True, BOSQUE, 0),
])

scope_data = [
    ("Venues", [
        "2–3 venue options per city",
        "Floorplan for each venue",
        "Equipment & furniture inventory",
        "Fixed pricing per format (meetup / workshop / hackathon)",
        "50 / 100 / 200-person configuration templates",
        "Google Calendar access (or 1h response backup — 1h text/email yes/no)",
        "iPad + badge printer as incentive for calendar access (500-unit pricing negotiation)",
        "Chatbot integration in venue website as additional contrepartie",
    ]),
    ("Catering", [
        "3–5 food options per city",
        "Low-cost: pizza & salad",
        "Middle: sandwich option",
        "High-end: elevated tacos or equivalent",
        "Additional: Asian food option",
        "Breakfast option",
        "Consistent menu across all cities for platform uniformity",
    ]),
    ("Media", [
        "1 photographer / videographer per city",
        "50 edited photos — raw folder included",
        "Up to 3-minute sizzle reel",
        "Consistent brief across all cities",
    ]),
]
cx = Inches(0.85)
for title, items in scope_data:
    add_textbox(s4, title, cx, Inches(2.32), Inches(4.0), Inches(0.28),
                font_name=CG, font_size=13, color=NEGRO)
    add_line(s4, cx, Inches(2.65), cx + Inches(3.8), Inches(2.65), color=BORDER)
    add_multiline(s4, cx, Inches(2.75), Inches(4.0), Inches(2.8), [
        (f"–  {it}", SG, 7.5, False, False, GRAY, 0 if i == 0 else 3)
        for i, it in enumerate(items)
    ])
    cx += Inches(4.22)

add_line(s4, Inches(0.85), Inches(5.72), Inches(12.5), Inches(5.72), color=BORDER)
add_rect(s4, Inches(0.85), Inches(5.84), Inches(0.04), Inches(0.55), fill_color=TIERRA_L)
add_textbox(s4,
    "All collected information will be compiled in a master spreadsheet with a linked folder of quotes. "
    "Each venue will receive a standardized floor plan template and equipment list per format configuration (50/100/200 pax).",
    Inches(1.05), Inches(5.82), Inches(11), Inches(0.65),
    font_name=CG, font_size=9.5, italic=True, color=GRAY_L)

footer(s4)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 — DELIVERABLES & WORKFLOW
# ═══════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK); bg(s5)
accent_line(s5); label(s5, "Deliverables & Workflow")

add_multiline(s5, Inches(0.85), Inches(1.2), Inches(10), Inches(1.0), [
    ("What you get.", CG, 30, False, False, NEGRO, 0),
    ("Ready to pitch.", CG, 30, False, True, BOSQUE, 0),
])

deliverables = [
    ("01", "Supplier Ecosystem", [
        "Master spreadsheet — all venues, caterers, media contacts per city",
        "Quote folder with all documentation received",
        "Floorplan templates for 50 / 100 / 200 person configurations",
    ]),
    ("02", "Digital Jungle Pilot", [
        "Pilot launch plan with Digital Jungle SF as anchor venue",
        "Branded materials and demo video structure",
        "Event format breakdown for each of the 3 formats",
    ]),
    ("03", "Investor-Ready Package", [
        "Fixed pricing catalog per city and format",
        "Proof of concept: venue partnerships concluded",
        "Google Calendar integration status per venue",
    ]),
]
cx = Inches(0.85)
cw = Inches(3.9)
for num, title, items in deliverables:
    add_textbox(s5, num, cx, Inches(2.45), cw, Inches(0.22),
                font_name=CG, font_size=10, color=BOSQUE)
    add_textbox(s5, title, cx, Inches(2.7), cw, Inches(0.38),
                font_name=CG, font_size=14, color=NEGRO)
    add_multiline(s5, cx, Inches(3.18), cw, Inches(2.5), [
        (f"–  {it}", SG, 8.5, False, False, GRAY, 0 if i == 0 else 5)
        for i, it in enumerate(items)
    ])
    cx += Inches(4.22)

add_line(s5, Inches(0.85), Inches(5.3), Inches(12.5), Inches(5.3), color=BORDER)
add_textbox(s5,
    "Sync cadence: daily or every 2 days alignment call throughout the 2-week sprint.",
    Inches(0.85), Inches(5.42), Inches(11.5), Inches(0.4),
    font_name=CG, font_size=10, italic=True, color=GRAY_L)

footer(s5)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 — INVESTMENT
# ═══════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK); bg(s6)
accent_line(s6); label(s6, "Investment")

add_multiline(s6, Inches(0.85), Inches(1.2), Inches(10), Inches(1.2), [
    ("Flat fee.", CG, 30, False, False, NEGRO, 0),
    ("Clean and simple.",  CG, 30, False, True,  BOSQUE, 4),
])

# Fee box
add_rect(s6, Inches(0.85), Inches(2.55), Inches(11.65), Inches(0.55), line_color=BORDER)
add_textbox(s6, "BLOOMING FLAT FEE", Inches(1.0), Inches(2.65), Inches(2.5), Inches(0.22),
            font_name=SG, font_size=6.5, color=BOSQUE)
add_textbox(s6, "$5,000", Inches(3.0), Inches(2.55), Inches(2.5), Inches(0.48),
            font_name=CG, font_size=26, color=NEGRO)
add_textbox(s6, "Upon invoice. Covers 2 weeks of full sourcing across all 10 cities — venues, catering, and media.",
            Inches(5.5), Inches(2.65), Inches(6.8), Inches(0.35),
            font_name=SG, font_size=8, color=GRAY_L)

# What's included
inc_y = Inches(3.32)
for title, items in [
    ("Included", [
        "2–3 venue options per city with full documentation",
        "3–5 catering options per city with fixed pricing",
        "1 photo/video contact per city with brief",
        "Master spreadsheet + quote folder",
        "Floor plan templates (50/100/200 pax)",
        "Pilot launch plan with Digital Jungle SF",
    ]),
    ("Ongoing (post-delivery)", [
        "20% commission on catalog pricing — per booking processed through EventGen",
        "Blooming manages supplier relationships as platform scales",
        "Additional cities sourced at same flat-fee model",
    ]),
]:
    pass

for cx, title, items in [
    (Inches(0.85), "Included in Flat Fee", [
        "2–3 venue options per city with full documentation",
        "3–5 catering options per city with fixed pricing",
        "1 photo/video contact per city with brief",
        "Master spreadsheet + complete quote folder",
        "Floor plan templates — 50 / 100 / 200 pax configs",
        "Pilot launch plan with Digital Jungle SF",
    ]),
    (Inches(7.0), "Ongoing Revenue Model", [
        "20% commission on catalog pricing per booking",
        "Blooming manages supplier relationships as platform scales",
        "Additional cities sourced at same flat-fee model",
        "Cadence: daily or every 2 days sync call during sprint",
    ]),
]:
    add_textbox(s6, title, cx, Inches(3.32), Inches(5.8), Inches(0.28),
                font_name=CG, font_size=13, italic=True, color=BOSQUE)
    add_multiline(s6, cx, Inches(3.65), Inches(5.8), Inches(2.2), [
        (f"–  {it}", SG, 8.5, False, False, GRAY, 0 if i == 0 else 4)
        for i, it in enumerate(items)
    ])

add_line(s6, Inches(0.85), Inches(5.72), Inches(12.5), Inches(5.72), color=BORDER)
add_textbox(s6,
    "Timeline: 2 weeks from May 16, 2026  ·  Deadline: May 30, 2026",
    Inches(0.85), Inches(5.84), Inches(11.5), Inches(0.3),
    font_name=CG, font_size=10, italic=True, color=GRAY_L)

footer(s6)

# ═══════════════════════════════════════════════════════════
# SLIDE 7 — CLOSE
# ═══════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK); bg(s7)

isotipo(s7, W/2, Inches(1.3), scale=1.1)

add_multiline(s7, Inches(1.5), Inches(1.85), Inches(10.3), Inches(1.9), [
    ("2 weeks.", CG, 38, False, False, NEGRO, 0),
    ("10 cities. Ready to pitch.", CG, 38, False, True, BOSQUE, 4),
], align=PP_ALIGN.CENTER)

add_textbox(s7, "Let's build the infrastructure.",
            Inches(1.5), Inches(3.72), Inches(10.3), Inches(0.32),
            font_name=CG, font_size=13, italic=True, color=GRAY_L, align=PP_ALIGN.CENTER)

add_textbox(s7,
    "EventGen needs a supplier ecosystem to pitch to investors with a real event engine running. "
    "Blooming has the relationships, the execution track record, and the continuity with Digital Jungle "
    "to deliver — in two weeks.",
    Inches(2.2), Inches(4.12), Inches(8.9), Inches(0.7),
    font_name=CG, font_size=11.5, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

steps = [
    "Sign and kick off — sprint starts from day of signature",
    "Blooming begins outreach to Digital Jungle SF as anchor pilot venue",
    "Parallel sourcing across all 10 cities — venues, catering, media",
    "Delivery by May 30: master spreadsheet + complete quote folder",
]
sy = Inches(4.95)
for step in steps:
    add_oval(s7, Inches(4.28), sy + Inches(0.08), Inches(0.07), Inches(0.07), BOSQUE)
    add_textbox(s7, step, Inches(4.45), sy, Inches(4.5), Inches(0.28),
                font_name=SG, font_size=8.5, color=GRAY)
    sy += Inches(0.34)

add_textbox(s7, "alison@weareblooming.co",
            Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.28),
            font_name=SG, font_size=10, color=GRAY, align=PP_ALIGN.CENTER)
add_textbox(s7, "WEAREBLOOMING.CO",
            Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.22),
            font_name=SG, font_size=7, color=BOSQUE, align=PP_ALIGN.CENTER)

footer(s7, "Blooming", "Growth that matters")

# ── Save ─────────────────────────────────────────────────
out = "proposals/clients/eventgen/Blooming x EventGen — Venues & Vendors Sourcing.pptx"
prs.save(out)
print(f"Saved → {out}")
