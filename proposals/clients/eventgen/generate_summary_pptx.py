from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NEGRO   = RGBColor(0x1A, 0x1A, 0x1A)
BLANCO  = RGBColor(0xFF, 0xFF, 0xFF)
BOSQUE  = RGBColor(0x5C, 0x6B, 0x5A)
TIERRA  = RGBColor(0x8B, 0x7D, 0x6B)
TIERRA_L= RGBColor(0xC4, 0xB8, 0xA8)
BORDER  = RGBColor(0xE0, 0xDD, 0xD8)
GRAY    = RGBColor(0x55, 0x55, 0x55)
GRAY_L  = RGBColor(0x88, 0x88, 0x88)
GRAY_F  = RGBColor(0xAA, 0xA8, 0xA4)

W = Inches(13.33)
H = Inches(7.5)
CG = "Cormorant Garamond"
SG = "Space Grotesk"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def bg(slide):
    add_rect(slide, 0, 0, W, H, fill_color=BLANCO)

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill_color:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color; shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_oval(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(9, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_line(slide, x1, y1, x2, y2, color=BORDER, width=Pt(0.5)):
    shape = slide.shapes.add_connector(1, x1, y1, x2, y2)
    shape.line.color.rgb = color; shape.line.width = width
    return shape

def add_textbox(slide, text, x, y, w, h, font_name=SG, font_size=11,
                bold=False, italic=False, color=NEGRO, align=PP_ALIGN.LEFT, word_wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = word_wrap
    tf = tb.text_frame; tf.word_wrap = word_wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name; run.font.size = Pt(font_size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_multiline(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, word_wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = word_wrap
    tf = tb.text_frame; tf.word_wrap = word_wrap
    first = True
    for (text, fname, fsize, fbold, fitalic, fcolor, spb) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if spb: p.space_before = Pt(spb)
        run = p.add_run()
        run.text = text
        run.font.name = fname; run.font.size = Pt(fsize)
        run.font.bold = fbold; run.font.italic = fitalic
        run.font.color.rgb = fcolor
    return tb

def accent_line(slide, x=Inches(0.85), y=Inches(1.08)):
    add_line(slide, x, y, x + Inches(0.32), y, color=NEGRO, width=Pt(0.75))

def label(slide, text, x=Inches(0.85), y=Inches(0.82)):
    add_textbox(slide, text.upper(), x, y, Inches(10), Inches(0.22),
                font_name=SG, font_size=7, color=BOSQUE)

def footer(slide, left="Blooming × EventGen", right="weareblooming.co"):
    y = H - Inches(0.48)
    add_textbox(slide, left,  Inches(0.85), y, Inches(5), Inches(0.28),
                font_name=SG, font_size=6, color=GRAY_F)
    add_textbox(slide, right, Inches(8.0),  y, Inches(5), Inches(0.28),
                font_name=SG, font_size=6, color=GRAY_F, align=PP_ALIGN.RIGHT)

def isotipo(slide, cx, cy, scale=1.0):
    s = scale
    add_oval(slide, cx - Inches(0.065*s), cy - Inches(0.26*s), Inches(0.13*s), Inches(0.32*s), BOSQUE)
    add_oval(slide, cx - Inches(0.065*s), cy - Inches(0.06*s), Inches(0.13*s), Inches(0.32*s), BOSQUE)
    add_oval(slide, cx - Inches(0.34*s),  cy - Inches(0.065*s), Inches(0.32*s), Inches(0.13*s), TIERRA)
    add_oval(slide, cx + Inches(0.02*s),  cy - Inches(0.065*s), Inches(0.32*s), Inches(0.13*s), TIERRA)
    r = Inches(0.065*s)
    add_oval(slide, cx - r, cy - r, r*2, r*2, NEGRO)


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK); bg(s1)

isotipo(s1, Inches(1.08), Inches(0.72), scale=0.9)
add_textbox(s1, "BLOOMING", Inches(1.38), Inches(0.62), Inches(3), Inches(0.28),
            font_name=SG, font_size=8, color=NEGRO)

add_multiline(s1, Inches(0.85), Inches(1.5), Inches(10), Inches(2.5), [
    ("BLOOMING × EventGen", CG, 42, False, False, NEGRO, 0),
    ("Summary of Actions.", CG, 42, False, True, BOSQUE, 4),
])

add_textbox(s1, "April 18 – June 5, 2026",
            Inches(0.85), Inches(3.85), Inches(9), Inches(0.32),
            font_name=SG, font_size=9, color=GRAY_L)

add_rect(s1, Inches(0.85), Inches(4.38), Inches(1.5), Inches(0.26), line_color=BORDER)
add_textbox(s1, "JUNE 2026", Inches(0.92), Inches(4.41), Inches(1.3), Inches(0.2),
            font_name=SG, font_size=6, color=GRAY_L)

add_textbox(s1, "Prepared by Alison Granger",
            Inches(0.85), Inches(6.5), Inches(7), Inches(0.28),
            font_name=SG, font_size=8, color=GRAY_L)

footer(s1, "Blooming", "weareblooming.co")


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — PIPELINE METRICS
# ═══════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK); bg(s2)
accent_line(s2); label(s2, "Pipeline Overview")

add_multiline(s2, Inches(0.85), Inches(1.2), Inches(11), Inches(1.0), [
    ("Pipeline Metrics Review.", CG, 30, False, False, NEGRO, 0),
    ("As of June 5, 2026.", CG, 30, False, True, BOSQUE, 4),
])

panels = [
    (Inches(0.85), "~90", "VENUES SOURCED", [
        ("~15", "CONFIRMED"),
        ("~18", "IN DISCUSSION"),
        ("~35", "EMAILED / OUTREACH"),
        ("~12", "FORM SENT"),
        ("~3",  "NOT INTERESTED"),
    ], "SF · LA · Austin · NY · Toronto · London · Paris · Berlin · Singapore — 9 cities"),
    (Inches(4.95), "~30", "MEDIA VENDORS SOURCED", [
        ("~18", "CONFIRMED"),
        ("~4",  "IN DISCUSSION"),
        ("~8",  "EMAILED"),
    ], "SF · LA · NY · Toronto · London · Paris · Berlin · Seattle · Austin · Singapore — 10 cities"),
    (Inches(9.05), "~120", "TOTAL CONTACTS", [
        ("~33",  "CONFIRMED RELATIONSHIPS"),
        ("~22",  "ACTIVE PIPELINE"),
        ("~27%", "CONVERSION RATE"),
    ], "Combined venues + media vendors"),
]

for cx, num, cat, breakdown, cities in panels:
    add_textbox(s2, num, cx, Inches(2.42), Inches(3.9), Inches(0.55),
                font_name=CG, font_size=30, color=NEGRO)
    add_textbox(s2, cat, cx, Inches(3.0), Inches(3.9), Inches(0.2),
                font_name=SG, font_size=6, color=BOSQUE)
    add_line(s2, cx, Inches(3.25), cx + Inches(3.8), Inches(3.25), color=BORDER)
    add_multiline(s2, cx, Inches(3.32), Inches(3.9), Inches(1.5), [
        (f"{v}  {l}", SG, 7.5, False, False, GRAY, 0 if i == 0 else 3)
        for i, (v, l) in enumerate(breakdown)
    ])
    add_textbox(s2, cities, cx, Inches(4.85), Inches(3.9), Inches(0.45),
                font_name=SG, font_size=6.5, color=GRAY_L)

add_line(s2, Inches(0.85), Inches(5.42), Inches(12.5), Inches(5.42), color=BORDER)

add_textbox(s2, "KEY HIGHLIGHTS", Inches(0.85), Inches(5.54), Inches(2), Inches(0.2),
            font_name=SG, font_size=6, color=BOSQUE)
add_multiline(s2, Inches(0.85), Inches(5.75), Inches(11.5), Inches(0.9), [
    (f"–  {h}", SG, 8, False, False, GRAY, 0 if i == 0 else 3)
    for i, h in enumerate([
        "Strong confirmation rate at 27%",
        "Berlin and Singapore fully sourced with confirmed vendors",
        "NY and LA have the deepest venue pipelines",
        "Media vendor network already operational in all 9 markets",
    ])
])

footer(s2)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — CITIES IN US
# ═══════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK); bg(s3)
accent_line(s3); label(s3, "Cities in US")

add_multiline(s3, Inches(0.85), Inches(1.2), Inches(8), Inches(0.8), [
    ("United States.", CG, 28, False, False, NEGRO, 0),
    ("Venue Pipeline.", CG, 28, False, True, BOSQUE, 4),
])

add_line(s3, Inches(0.85), Inches(2.08), Inches(12.5), Inches(2.08), color=BORDER)
add_textbox(s3, "CITY", Inches(0.85), Inches(2.16), Inches(1.85), Inches(0.2),
            font_name=SG, font_size=6, color=BOSQUE)
add_textbox(s3, "VENUES & STATUS", Inches(2.9), Inches(2.16), Inches(9.8), Inches(0.2),
            font_name=SG, font_size=6, color=BOSQUE)
add_line(s3, Inches(0.85), Inches(2.38), Inches(12.5), Inches(2.38), color=BORDER)

CITY_X  = Inches(0.85)
VENUE_X = Inches(2.9)
VENUE_W = Inches(9.8)

us_cities = [
    ("San Francisco", [
        "Digital Garage — need to work on layout, partnership on",
    ]),
    ("New York", [
        "The Box House Hotel",
        "VESPER",
        "⇒ Call Tuesday to streamline details — catering in-house, AV: TBC (not confirmed if included)",
    ]),
    ("Los Angeles", [
        "Hype Studios — great venue",
        "Hot Shot Muffler — too small but great AV",
        "Skyline Penthouse — waiting for specs",
    ]),
    ("Austin", [
        "Wanderlust — Shady Lane",
        "Wanderlust — Downtown",
        "Wanderlust — Barton Springs",
        "Regal Room",
        "dadaLab — call on Wednesday",
    ]),
    ("Seattle", [
        "1712 Studios",
        "Block 41 — need to send the layout",
    ]),
]

cy = Inches(2.48)
for city, venues in us_cities:
    row_h = Inches(max(len(venues) * 0.19 + 0.06, 0.25))
    add_textbox(s3, city, CITY_X, cy, Inches(1.85), row_h,
                font_name=CG, font_size=10, color=BOSQUE)
    add_multiline(s3, VENUE_X, cy, VENUE_W, row_h, [
        (v, SG, 7.5, False, False, GRAY, 0 if i == 0 else 2)
        for i, v in enumerate(venues)
    ])
    cy += row_h + Inches(0.1)
    add_line(s3, CITY_X, cy, Inches(12.5), cy, color=BORDER, width=Pt(0.3))
    cy += Inches(0.09)

footer(s3)


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — CITIES OUT OF US
# ═══════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK); bg(s4)
accent_line(s4); label(s4, "Cities out of US")

add_multiline(s4, Inches(0.85), Inches(1.2), Inches(8), Inches(0.8), [
    ("International.", CG, 28, False, False, NEGRO, 0),
    ("Venue Pipeline.", CG, 28, False, True, BOSQUE, 4),
])

add_line(s4, Inches(0.85), Inches(2.08), Inches(12.5), Inches(2.08), color=BORDER)
add_textbox(s4, "CITY", Inches(0.85), Inches(2.16), Inches(1.85), Inches(0.2),
            font_name=SG, font_size=6, color=BOSQUE)
add_textbox(s4, "VENUES & STATUS", Inches(2.9), Inches(2.16), Inches(9.8), Inches(0.2),
            font_name=SG, font_size=6, color=BOSQUE)
add_line(s4, Inches(0.85), Inches(2.38), Inches(12.5), Inches(2.38), color=BORDER)

intl_cities = [
    ("London", [
        "Tobacco Dock — complex venue (must work with their caterer). Call with sales director June 5th — want to work with us but several rounds needed. Commission: 8%",
        "New venue recommendations not yet contacted (from vendors): acceleratorcentre.com",
    ]),
    ("Paris", [
        "0 fully confirmed (150+ pax) — warm lead with Station F, coming back week of June 8th",
        "Equiria — 50 pax in theater",
        "mk2 — call on Wednesday",
        "La Maison (lamaison.ai) — need to circle back",
    ]),
    ("Berlin", [
        "SpaceShack — confirmed partner, very happy to work with us",
        "Forum Factory + Backspace Berlin — same entity",
    ]),
    ("Toronto", [
        "Startuptive — 75 pax in theater style (too small)",
        "CSI — warm lead, busy this week, coming back week of June 8th",
    ]),
    ("Singapore", [
        "1 confirmed (small venue)",
        "1 in conversation — call on Wednesday",
        "Need to reactivate additional contact",
    ]),
]

cy = Inches(2.48)
for city, venues in intl_cities:
    row_h = Inches(max(len(venues) * 0.2 + 0.06, 0.25))
    add_textbox(s4, city, CITY_X, cy, Inches(1.85), row_h,
                font_name=CG, font_size=10, color=BOSQUE)
    add_multiline(s4, VENUE_X, cy, VENUE_W, row_h, [
        (v, SG, 7.5, False, False, GRAY, 0 if i == 0 else 2)
        for i, v in enumerate(venues)
    ])
    cy += row_h + Inches(0.1)
    add_line(s4, CITY_X, cy, Inches(12.5), cy, color=BORDER, width=Pt(0.3))
    cy += Inches(0.09)

footer(s4)


# ═══════════════════════════════════════════════════════════
# SLIDE 5 — CATERERS
# ═══════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK); bg(s5)
accent_line(s5); label(s5, "Caterers")

add_multiline(s5, Inches(0.85), Inches(1.2), Inches(10), Inches(1.0), [
    ("Catering.", CG, 30, False, False, NEGRO, 0),
    ("Sourcing pipeline.", CG, 30, False, True, BOSQUE, 4),
])

add_rect(s5, Inches(0.85), Inches(2.55), Inches(11.65), Inches(0.55), line_color=BORDER)
for cx, num, cat in [
    (Inches(1.2),  "~45", "TOTAL CONTACTS"),
    (Inches(3.75), "~27", "EMAILED"),
    (Inches(6.3),  "~15", "FORM SENT"),
    (Inches(8.85), "~3",  "CONFIRMED"),
]:
    add_textbox(s5, num, cx, Inches(2.56), Inches(2.3), Inches(0.38),
                font_name=CG, font_size=22, color=NEGRO)
    add_textbox(s5, cat, cx, Inches(2.94), Inches(2.3), Inches(0.18),
                font_name=SG, font_size=5.5, color=BOSQUE)

add_line(s5, Inches(0.85), Inches(3.22), Inches(12.5), Inches(3.22), color=BORDER)

for cx, title, items in [
    (Inches(0.85), "Confirmed", [
        "Los Angeles — MBM Hospitality",
        "Los Angeles — Pino's Sandwiches",
        "Berlin — Bellabona",
    ]),
    (Inches(5.5), "In-House (venue-provided)", [
        "London (Tobacco Dock)",
        "New York (Box House Hotel / VESPER)",
    ]),
    (Inches(9.5), "Next Pipeline", [
        "Austin — recommendations from Regal Room pending",
    ]),
]:
    add_textbox(s5, title, cx, Inches(3.35), Inches(3.5), Inches(0.28),
                font_name=CG, font_size=13, italic=True, color=BOSQUE)
    add_multiline(s5, cx, Inches(3.68), Inches(3.5), Inches(2.0), [
        (f"–  {it}", SG, 8.5, False, False, GRAY, 0 if i == 0 else 4)
        for i, it in enumerate(items)
    ])

add_line(s5, Inches(0.85), Inches(5.6), Inches(12.5), Inches(5.6), color=BORDER)
add_rect(s5, Inches(0.85), Inches(5.72), Inches(0.04), Inches(0.55), fill_color=TIERRA_L)
add_textbox(s5,
    "Catering has been harder to source — requires catching contacts at the right time. "
    "In-house catering at London and NYC venues simplifies logistics for those markets.",
    Inches(1.05), Inches(5.70), Inches(11), Inches(0.65),
    font_name=CG, font_size=9.5, italic=True, color=GRAY_L)

footer(s5)


# ═══════════════════════════════════════════════════════════
# SLIDE 6 — MEDIA VENDORS
# ═══════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK); bg(s6)
accent_line(s6); label(s6, "Media Vendors")

add_multiline(s6, Inches(0.85), Inches(1.2), Inches(10), Inches(1.0), [
    ("Media.", CG, 30, False, False, NEGRO, 0),
    ("Photo & Video sourcing.", CG, 30, False, True, BOSQUE, 4),
])

for cx, num, cat in [
    (Inches(0.85), "~30", "TOTAL VENDORS"),
    (Inches(4.0),  "~18", "CONFIRMED"),
    (Inches(7.2),  "~4",  "IN DISCUSSION"),
    (Inches(10.4), "~8",  "EMAILED"),
]:
    add_textbox(s6, num, cx, Inches(2.45), Inches(3.0), Inches(0.52),
                font_name=CG, font_size=30, color=NEGRO)
    add_textbox(s6, cat, cx, Inches(3.0),  Inches(3.0), Inches(0.22),
                font_name=SG, font_size=6,  color=BOSQUE)

add_line(s6, Inches(0.85), Inches(3.32), Inches(12.5), Inches(3.32), color=BORDER)

add_textbox(s6, "CITIES COVERED", Inches(0.85), Inches(3.45), Inches(2.5), Inches(0.2),
            font_name=SG, font_size=6, color=BOSQUE)
add_textbox(s6,
    "San Francisco  ·  Los Angeles  ·  Paris  ·  New York  ·  Toronto  ·  London  ·  Berlin  ·  Seattle  ·  Austin  ·  Singapore",
    Inches(0.85), Inches(3.68), Inches(11.5), Inches(0.32),
    font_name=CG, font_size=13, color=NEGRO)
add_textbox(s6, "10 cities — full media vendor coverage",
            Inches(0.85), Inches(4.05), Inches(11.5), Inches(0.28),
            font_name=SG, font_size=8, color=GRAY_L)

add_line(s6, Inches(0.85), Inches(4.45), Inches(12.5), Inches(4.45), color=BORDER)

add_rect(s6, Inches(0.85), Inches(4.55), Inches(0.04), Inches(0.6), fill_color=TIERRA_L)
add_textbox(s6,
    "Media vendor network is already operational across all 10 markets. "
    "~60% confirmation rate across sourced contacts — the strongest track in the overall pipeline.",
    Inches(1.05), Inches(4.55), Inches(11), Inches(0.65),
    font_name=CG, font_size=10, italic=True, color=GRAY_L)

footer(s6, "Blooming × EventGen", "weareblooming.co")


# ── Save ────────────────────────────────────────────────────
out = "proposals/clients/eventgen/Blooming x EventGen — Summary of Actions.pptx"
prs.save(out)
print(f"Saved → {out}")
